"""This script plots histograms created in the analysis of Jetscape events

.. codeauthor:: James Mulligan <james.mulligan@berkeley.edu>, UC Berkeley
.. codeauthor:: Raymond Ehlers <raymond.ehlers@cern.ch>, LBL/UCB
"""

from __future__ import annotations

# General
import argparse
import ctypes
import logging
import sys
from pathlib import Path
from typing import Any

import h5py
import numpy as np
import pandas as pd
import ROOT  # pyright: ignore [reportMissingImports]
import uproot
import yaml
from jetscape_analysis.base import common_base, helpers
from jetscape_analysis.data_curation import observable as observable_module
from plot import plot_results_STAT_utils

# Prevent ROOT from stealing focus when plotting
ROOT.gROOT.SetBatch(True)

logger = logging.getLogger(__name__)

_model_display_name = {
    "jetscape": "JETSCAPE",
    "hybrid": "Hybrid",
}


################################################################
class PlotResults(common_base.CommonBase):
    # ---------------------------------------------------------------
    # Constructor
    # ---------------------------------------------------------------
    def __init__(
        self,
        config_file: str | Path = "",
        input_file: str | Path = "",
        output_dir: str | Path | None = "",
        pp_ref_file: str | Path = "",
        model_name: str = "jetscape",
        **kwargs,
    ):
        super().__init__(**kwargs)

        # Validation
        config_file = Path(config_file)
        self.config_file = config_file
        input_file = Path(input_file)
        pp_ref_file = Path(pp_ref_file)

        # Encoder context for the migrated jet/substructure observables (set inside the migrated
        # branch of plot_jet_observables, read by get_histogram). Default OFF so every non-jet
        # get_histogram call (hadron, semi-inclusive, STAR, ...) keeps the legacy f-string name.
        self._is_migrated_obs = False
        self._encoder_block = None
        self._encoder_grooming_setting = None
        self._encoder_pt_bin = None
        self._encoder_axis_entry = None
        self._encoder_charge = None
        self._encoder_angularity = None

        # Default to the same directory as the input_file if an output_dir is not provided.
        if not output_dir:
            output_dir = input_file.parent
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True, parents=True)

        self.plot_utils = plot_results_STAT_utils.PlotUtils()
        self.plot_utils.setOptions()
        ROOT.gROOT.ForceStyle()

        self.input_file = ROOT.TFile(str(input_file), "READ")

        self.data_color = ROOT.kGray + 3
        self.data_marker = 21
        self.jetscape_color = [
            # RJE changed kViolet to kGreen to enhance the clarity in the plots (e.g. ensure we can see when bands are overlapping).
            # n.b. the new colors may be worse for colorblindness, so I kept the original colors below.
            # ROOT.kViolet - 8,
            ROOT.kGreen - 6,
            ROOT.kViolet - 8,
            ROOT.kRed - 7,
            ROOT.kTeal - 8,
            ROOT.kCyan - 2,
            ROOT.kGreen - 6,
            ROOT.kAzure - 4,
            ROOT.kOrange + 6,
            ROOT.kBlue - 10,
        ]
        self.jetscape_fillstyle = [1001, 3144, 1001, 3144]
        self.jetscape_alpha = [0.7, 0.7, 0.7, 0.7]
        self.jetscape_marker = 20
        self.marker_size = 1.5
        self.line_width = 2
        self.line_style = 1
        self.file_format = ".pdf"

        # Check whether pp or AA
        self.is_AA = False
        # NOTE: We only use this in AA, but it's better for it to always be defined
        self.observable_centrality_list = []
        if "PbPb" in str(input_file) or "AuAu" in str(input_file):
            self.is_AA = True

        # If AA, load the pp reference results so that we can construct RAA
        if self.is_AA:
            self.pp_ref_file = ROOT.TFile(str(pp_ref_file), "READ")

        # Store the model name for customization
        if model_name == "":
            model_name = "jetscape"
            msg = "No model name provided, so default to jetscape."
            logger.warning(msg)
        self.model_name = model_name
        if self.model_name not in _model_display_name:
            msg = f"Provided unexpected model name {self.model_name} that is not supported. Please check"
            raise ValueError(msg)

        # Read config file
        with config_file.open() as stream:
            self.config = yaml.safe_load(stream)
        self.sqrts = self.config["sqrt_s"]
        self.power = self.config["power"]
        self.pt_ref = self.config["pt_ref"]

        # Observable objects keyed "{sqrts}_{observable_class}_{name}", used to rebuild the
        # encoder-based histogram names for the migrated groomed & substructure observables so the
        # plotter finds what the histogrammer wrote (see MIGRATION_NOTES Step 4 ④, OBSERVABLE_EDGE_CASES
        # C5). Container-safe import (attrs/numpy/pyyaml only), same as the histogrammer.
        self.observables_info = observable_module.read_observables_from_config(
            config_path=self.config_file, sqrt_s=self.sqrts
        )

        # If AA, set different options for hole subtraction treatment
        self.jet_collection_labels_AA = self.config["jet_collection_labels"]  # + ['_shower_recoil_unsubtracted']
        self.jet_collection_label_pp = ""
        if self.is_AA:
            self.jet_collection_labels = self.jet_collection_labels_AA
        else:
            self.jet_collection_labels = [self.jet_collection_label_pp]

        # We will write final results after all scaling steps, along with data, to file
        self.output_dict = {}

        logger.info(self)

    # -------------------------------------------------------------------------------------------
    # -------------------------------------------------------------------------------------------
    # -------------------------------------------------------------------------------------------
    def plot_results(self):
        self.analysis = "Analysis1"

        if self.analysis == "hadron_jet_RAA":
            self.plot_hadron_observables(observable_type="hadron")

            self.plot_jet_observables(observable_type="inclusive_chjet")

            if "inclusive_jet" in self.config:
                self.plot_jet_observables(observable_type="inclusive_jet")

        else:
            self.plot_hadron_observables(observable_type="hadron")

            self.plot_jet_observables(observable_type="inclusive_chjet")

            if "inclusive_jet" in self.config:
                self.plot_jet_observables(observable_type="inclusive_jet")

            if "dijet_trigger_jet" in self.config:
                self.plot_jet_observables(observable_type="dijet_trigger_jet")

            # Semi-inclusive hadron+jet recoil observables (IAA_pt_alice, dphi_alice).
            # The method existed but was never invoked here, so these never plotted.
            if "hadron_trigger_chjet" in self.config:
                self.plot_hadron_trigger_chjet_observables(observable_type="hadron_trigger_chjet")

        self.plot_event_qa()

        self.write_output_objects()

        # Generate pptx for convenience
        if self.file_format == ".png":
            self.generate_pptx()

    # -------------------------------------------------------------------------------------------
    # Plot hadron observables
    # -------------------------------------------------------------------------------------------
    def plot_hadron_observables(self, observable_type=""):
        logger.info(f"\nPlot {observable_type} observables...")

        for observable, block in self.config[observable_type].items():
            for centrality_index, centrality in enumerate(block["centrality"]):
                if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                    continue

                # Initialize observable configuration
                self.suffix = ""
                self.init_observable(observable_type, observable, block, centrality, centrality_index)

                # Plot observable
                self.plot_observable(observable_type, observable, centrality)

    # -------------------------------------------------------------------------------------------
    # Plot hadron correlation observables
    # -------------------------------------------------------------------------------------------
    def plot_hadron_correlation_observables(self, observable_type: str) -> None:
        logger.info(f"\nPlot {observable_type} observables...")

        for observable, block in self.config[observable_type].items():
            for centrality_index, centrality in enumerate(block["centrality"]):
                if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                    continue

                # for hadron v2
                if "v2" in observable:
                    self.init_observable(observable_type, observable, block, centrality, centrality_index)
                    # Histogram observable
                    self.plot_observable(observable_type, observable, centrality)

    # -------------------------------------------------------------------------------------------
    # Plot hadron triggered hadron observables
    # -------------------------------------------------------------------------------------------
    def plot_hadron_trigger_hadron_observables(self, observable_type: str) -> None:
        logger.info(f"\nPlot {observable_type} observables...")

        for observable, block in self.config[observable_type].items():
            for centrality_index, centrality in enumerate(block["centrality"]):
                if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                    continue

                # STAR dihadron
                if observable == "dihadron_star":
                    # Initialize observable configuration
                    pt_trigger_ranges = block["pt_trig"]
                    pt_associated_ranges = block["pt_assoc"]
                    # Loop over trigger and associated ranges
                    for pt_trig_min, pt_trig_max in pt_trigger_ranges:
                        for pt_assoc_min, pt_assoc_max in pt_associated_ranges:
                            # If the upper range has -1, it's unbounded, so we make it large enough not to matter
                            pt_assoc_max = 1000 if pt_assoc_max == -1 else pt_assoc_max  # noqa: PLW2901
                            self.suffix = (
                                f"_pt_trig_{pt_trig_min:g}_{pt_trig_max:g}_pt_assoc_{pt_assoc_min:g}_{pt_assoc_max:g}"
                            )
                            self.init_observable(observable_type, observable, block, centrality, centrality_index)

                            # Histogram observable
                            self.plot_observable(observable_type, observable, centrality)

    # -------------------------------------------------------------------------------------------
    # Histogram inclusive jet observables
    # -------------------------------------------------------------------------------------------
    def plot_jet_observables(self, observable_type: str = "") -> None:  # noqa: C901
        logger.info(f"\nPlot {observable_type} observables...")

        for observable, block in self.config[observable_type].items():
            if not block.get("enabled", True):
                continue
            for centrality_index, centrality in enumerate(block["centrality"]):
                # Custom skip to only process inclusive hadron and jet RAA.
                # Used historically - not really needed in July 2025, but we leave it as is.
                if self.analysis == "hadron_jet_RAA" and not observable.startswith("pt_"):
                    logger.info(f"skipping {observable}")
                    continue

                for self.jet_R in block["jet"]["R"]:
                    # Optional: Loop through pt bins. xj_atlas (dijet) is binned by the LEADING/trigger
                    # jet pt (mirrors the histogrammer); only 0-10% reports all leading-pt bins.
                    pt_edges = block["trigger"]["pt"] if observable == "xj_atlas" else block["jet"]["pt"]
                    for pt_bin in range(len(pt_edges) - 1):
                        if observable == "xj_atlas" and centrality_index > 0 and pt_bin != 0:
                            continue
                        pt_suffix = ""
                        if len(pt_edges) > 2:
                            pt_suffix = f"_pt{pt_bin}"

                        # Optional: subobservable. Each entry is (label, axis_entry, charge_value,
                        # angularity_spec), mirroring the histogrammer: the label feeds the (jet_R/pt-only)
                        # suffix used for plot labels/filenames; axis_entry feeds the encoder for migrated
                        # axis observables (the bare `type` label can't distinguish the grooming variant);
                        # charge_value (kappa) feeds the encoder for charge_cms; angularity_spec (alpha)
                        # feeds it for the angularity observables. Keys are mutually exclusive in the YAML.
                        # Tuple is (subobservable_label, axis_entry, charge_value, angularity_spec,
                        # subjet_spec). Step 6 added the 5th slot for zr_alice's subjet_R.
                        subobservable_label_list = [("", None, None, None, None)]
                        if "axis" in block["jet"]:
                            subobservable_label_list = [
                                (f"_{axis_block['type']}", axis_block, None, None, None) for axis_block in block["jet"]["axis"]
                            ]
                        if "charge" in block["jet"]:
                            subobservable_label_list = [
                                (f"_k{kappa}", None, kappa, None, None) for kappa in block["jet"]["charge"]
                            ]
                        # Step 5: angularity (alpha) sub-observable -- one AngularitySpec per alpha, feeding
                        # the encoder name + per-alpha data-block overlay table selection (jet_angularity).
                        if "angularity" in block["jet"]:
                            subobservable_label_list = [
                                (f"_alpha{a['alpha']}", None, None, observable_module.AngularitySpec(alpha=a["alpha"]), None)
                                for a in block["jet"]["angularity"]
                            ]
                        # Step 6: subjet_R (zr_alice) -- one SubjetRSpec per subjet R, feeding the `_r{r}`
                        # legacy suffix (matching the analyzer column) + per-r data-block overlay table
                        # (jet_subjet_R key). zr_alice stays on the legacy name path (NOT the encoder).
                        if "subjet_R" in block["jet"]:
                            subobservable_label_list = [
                                (f"_r{r}", None, None, None, observable_module.SubjetRSpec(r=r))
                                for r in block["jet"]["subjet_R"]
                            ]

                        # Whether the histogram name comes from the observable encoder (the analyzer +
                        # histogrammer use encoder names for these) or the legacy hand-built f-string.
                        is_migrated = (
                            observable_type,
                            observable,
                        ) in plot_results_STAT_utils.ENCODER_MIGRATED_JET_OBSERVABLES
                        # jet_axis is an essential (encoded) parameter only when >1 axis variant is
                        # configured; with a single axis entry the encoder omits it (mirrors the
                        # histogrammer: axis_alice passes jet_axis, axis_cms does not).
                        axis_is_essential = isinstance(block["jet"].get("axis"), list) and len(block["jet"]["axis"]) > 1

                        for subobservable_label, axis_entry, charge_value, angularity_spec, subjet_spec in subobservable_label_list:
                            # Set normalization
                            self_normalize = False
                            # "axis" -> jet-axis-difference observables (axis_alice/axis_cms) report a
                            # self-normalized 1/sigma_inc dN/dDeltaR distribution, so the MC must be
                            # self-normalized too (without it the raw MC sits ~0 against the data scale).
                            # "zr" -> zr_alice's leading-subjet z_r is also self-normalized (1/sigma).
                            for x in ["mass", "g", "ptd", "charge", "mg", "zg", "tg", "ktg", "xj", "axis", "zr"]:
                                if x in observable:
                                    self_normalize = True

                            if "grooming_settings" in block["jet"]:
                                for grooming_setting in block["jet"]["grooming_settings"]:
                                    # New YAML schema mixes grooming methods (soft_drop + dynamical_grooming).
                                    # Dynamical Grooming (DyG) is supported only on the encoder path
                                    # (migrated observables); non-soft-drop on a non-migrated observable has
                                    # no legacy name -> skip. Kept in lockstep with the histogrammer loop.
                                    grooming_type = grooming_setting.get("type", "soft_drop")
                                    if grooming_type != "soft_drop" and not is_migrated:
                                        continue
                                    # Custom skip
                                    if observable in ["zg_alice", "tg_alice"]:
                                        if np.isclose(self.jet_R, 0.4) and centrality_index == 0:
                                            continue
                                        if np.isclose(self.jet_R, 0.2) and centrality_index == 1:
                                            continue

                                    logger.info(f"      grooming_setting = {grooming_setting}")
                                    if grooming_type == "soft_drop":
                                        zcut = grooming_setting["z_cut"]
                                        beta = grooming_setting["beta"]
                                        # Option to take zcut and beta = 0 as the ungroomed case, where we fall back to the standard suffix
                                        if np.isclose(zcut, 0.0) and np.isclose(beta, 0.0):
                                            self.suffix = f"_R{self.jet_R}{subobservable_label}"
                                        else:
                                            self.suffix = f"_R{self.jet_R}_zcut{zcut}_beta{beta}{subobservable_label}"
                                    else:
                                        # Dynamical grooming (migrated path only, gated above)
                                        self.suffix = f"_R{self.jet_R}_a{grooming_setting['a']}{subobservable_label}"

                                    if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                                        continue

                                    # For migrated groomed observables, select the per-grooming
                                    # HEPData overlay table (data block keyed by jet_grooming_settings
                                    # in the GroomingSettingsSpec encoding, WITH the SD_/DyG_ prefix --
                                    # distinct from the no-prefix column encoding used for the name).
                                    data_block_params = None
                                    if is_migrated:
                                        data_block_params = {
                                            "jet_grooming_settings": observable_module.GroomingSettingsSpec(
                                                grooming_setting
                                            ).encode()
                                        }
                                        # Step 5: per-alpha overlay table selection for angularity observables.
                                        if angularity_spec is not None:
                                            data_block_params["jet_angularity"] = angularity_spec.encode()
                                    # Stash the encoder context so get_histogram rebuilds the migrated
                                    # histogram name from the same encoder the histogrammer used.
                                    self._is_migrated_obs = is_migrated
                                    self._encoder_block = block
                                    self._encoder_grooming_setting = grooming_setting
                                    self._encoder_pt_bin = pt_bin
                                    self._encoder_axis_entry = None
                                    self._encoder_charge = None
                                    self._encoder_angularity = angularity_spec

                                    # Initialize observable configuration
                                    self.init_observable(
                                        observable_type,
                                        observable,
                                        block,
                                        centrality,
                                        centrality_index,
                                        pt_suffix=pt_suffix,
                                        self_normalize=self_normalize,
                                        data_block_params=data_block_params,
                                    )

                                    # Plot observable
                                    self.plot_observable(observable_type, observable, centrality, pt_suffix)
                                    self._is_migrated_obs = False

                            else:
                                self.suffix = f"_R{self.jet_R}{subobservable_label}"
                                if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                                    continue

                                # For migrated axis observables with >1 axis variant (axis_alice), the
                                # data block has one HEPData overlay table per jet_axis variant, so
                                # select it (JetAxisDifferenceSpec encoding -- grooming carried inside
                                # the axis). NOTE: axis_alice still carries legacy hepdata_* keys, so
                                # its overlay routes through tgraph_from_hepdata (ROOT) and ignores
                                # this; it takes effect for any data:-only axis observable (Step 6).
                                data_block_params = None
                                if is_migrated and axis_is_essential and axis_entry is not None:
                                    data_block_params = {
                                        "jet_axis": observable_module.JetAxisDifferenceSpec(
                                            type=axis_entry["type"], grooming_settings=axis_entry.get("grooming_settings")
                                        ).encode()
                                    }
                                # charge_cms: select the per-kappa HEPData overlay table (data block
                                # keyed by jet_charge in the JetChargeSpec encoding, kappa_X).
                                if is_migrated and charge_value is not None:
                                    data_block_params = {
                                        "jet_charge": observable_module.JetChargeSpec(charge_value).encode()
                                    }
                                # Step 6: zr_alice -- one HEPData overlay table per subjet R; select it by
                                # the jet_subjet_R key (SubjetRSpec encoding, "r_{r}"). NOT gated on
                                # is_migrated (zr_alice is on the legacy name path, but its overlay still
                                # resolves through the per-r data-block table).
                                if subjet_spec is not None:
                                    data_block_params = {"jet_subjet_R": subjet_spec.encode()}
                                # Stash the encoder context (no grooming in the else branch; pass the
                                # axis entry only when jet_axis is essential, mirroring the analyzer;
                                # charge_value is the kappa for charge_cms).
                                self._is_migrated_obs = is_migrated
                                self._encoder_block = block
                                self._encoder_grooming_setting = None
                                self._encoder_pt_bin = pt_bin
                                self._encoder_axis_entry = axis_entry if axis_is_essential else None
                                self._encoder_charge = charge_value
                                self._encoder_angularity = angularity_spec

                                # Initialize observable configuration
                                self.init_observable(
                                    observable_type,
                                    observable,
                                    block,
                                    centrality,
                                    centrality_index,
                                    pt_suffix=pt_suffix,
                                    self_normalize=self_normalize,
                                    data_block_params=data_block_params,
                                )

                                # Plot observable
                                self.plot_observable(observable_type, observable, centrality, pt_suffix)
                                self._is_migrated_obs = False

    # -------------------------------------------------------------------------------------------
    # Histogram semi-inclusive jet observables
    # -------------------------------------------------------------------------------------------
    def plot_hadron_trigger_chjet_observables(self, observable_type=""):
        logger.info(f"\nPlot {observable_type} observables...")

        for observable, block in self.config[observable_type].items():
            if not block.get("enabled", True):
                continue
            for centrality_index, centrality in enumerate(block["centrality"]):
                for self.jet_R in block["jet"]["R"]:
                    self.suffix = f"_R{self.jet_R}"
                    if "hepdata" not in block and "custom_data" not in block and "data" not in block:
                        continue

                    # Set normalization
                    self_normalize = False
                    for x in ["nsubjettiness"]:
                        if x in observable:
                            self_normalize = True

                    # dphi is differential in recoil-jet pt (one _pt{i} per jet-pt bin -> a separate
                    # Delta_recoil(Dphi) distribution and overlay table); IAA is pt-integrated.
                    if "dphi" in observable and isinstance(block.get("jet", {}).get("pt"), list):
                        pt_suffixes = [f"_pt{i}" for i in range(len(block["jet"]["pt"]) - 1)]
                    else:
                        pt_suffixes = [""]

                    for pt_suffix in pt_suffixes:
                        # Initialize observable configuration
                        self.init_observable(
                            observable_type,
                            observable,
                            block,
                            centrality,
                            centrality_index,
                            pt_suffix=pt_suffix,
                            self_normalize=self_normalize,
                        )

                        # Plot observable
                        self.plot_observable(observable_type, observable, centrality, pt_suffix=pt_suffix)

    # -------------------------------------------------------------------------------------------
    # Initialize a single observable's config
    # -------------------------------------------------------------------------------------------
    def init_observable(
        self,
        observable_type,
        observable,
        block,
        centrality,
        centrality_index,
        pt_suffix="",
        self_normalize=False,
        data_block_params=None,
    ):
        # Initialize an empty dict containing relevant info
        self.observable_settings = {}

        # Initialize common settings into class members
        self.init_common_settings(observable, block)

        # -----------------------------------------------------------
        # Initialize data distribution into self.observable_settings
        _available_hepdata_files = plot_results_STAT_utils.available_hepdata_files_in_block(block)
        if any(_available_hepdata_files):
            self.observable_settings["data_distribution"] = self.plot_utils.tgraph_from_hepdata(
                block,
                self.is_AA,
                self.sqrts,
                observable_type,
                observable,
                centrality_index,
                suffix=self.suffix,
                pt_suffix=pt_suffix,
            )
        elif "custom_data" in block:
            self.observable_settings["data_distribution"] = self.plot_utils.tgraph_from_yaml(
                block,
                self.is_AA,
                self.sqrts,
                observable_type,
                observable,
                centrality_index,
                suffix=self.suffix,
                pt_suffix=pt_suffix,
            )
        elif "data" in block and not (self.skip_pp and not self.is_AA):
            # New-schema `data:` block -> read the measured points from the HEPData YAML.
            # Note: pass `centrality` (the [low, high] list), not centrality_index — the
            # exact table+index is matched by centrality/jet_R/jet_pt value.
            # skip_pp marks an observable with no pp measurement (e.g. pt_y_atlas, a PbPb
            # |y| self-ratio whose only |y| HEPData table is PbPb). For the pp arm we then
            # load NO data overlay -- otherwise an off-axis/unrelated table (the pp pT
            # spectrum) would draw a misleading "Data" legend contradicting the
            # "skip data plot -- no pp data in HEPData" annotation.
            self.observable_settings["data_distribution"] = self.plot_utils.tgraph_from_data_block(
                block,
                self.is_AA,
                self.sqrts,
                observable_type,
                observable,
                centrality,
                suffix=self.suffix,
                pt_suffix=pt_suffix,
                data_block_params=data_block_params,
            )
        else:
            self.observable_settings["data_distribution"] = None

        # Also initialize systematic uncertainties
        # logger.info(block)

        # -----------------------------------------------------------
        # Initialize JETSCAPE distribution into self.observable_settings
        self.initialize_jetscape_distribution(
            observable_type,
            observable,
            block,
            centrality,
            centrality_index,
            pt_suffix=pt_suffix,
            self_normalize=self_normalize,
        )

        # -----------------------------------------------------------
        # For pp case -- form ratio of JETSCAPE to data and load into self.observable_settings
        if not self.is_AA:
            if (
                self.observable_settings["data_distribution"]
                and self.observable_settings["jetscape_distribution"]
                and not self.observable_settings["jetscape_distribution"].InheritsFrom(ROOT.TH2.Class())
                and not self.skip_pp_ratio
            ):
                self.observable_settings["ratio"] = self.plot_utils.divide_histogram_by_tgraph(
                    self.observable_settings["jetscape_distribution"],
                    self.observable_settings["data_distribution"],
                    include_tgraph_uncertainties=False,
                )
            else:
                self.observable_settings["ratio"] = None
        if "v2" in observable and self.is_AA and self.observable_settings["jetscape_distribution"]:
            self.observable_settings["ratio"] = self.plot_utils.divide_histogram_by_tgraph(
                self.observable_settings["jetscape_distribution"], self.observable_settings["data_distribution"]
            )

    # -------------------------------------------------------------------------------------------
    # Initialize from settings from config file into class members
    # -------------------------------------------------------------------------------------------
    def init_common_settings(self, observable: str, block: dict[str, Any]) -> None:  # noqa: C901
        self.xtitle = block.get("xtitle", "")
        # Acceptance cuts. eta (pseudorapidity) -> self.eta_cut; rapidity (y) -> self.y_cut; these
        # are distinct and drive different scale_histogram normalizations (e.g. ATLAS jets use
        # rapidity, CMS jets/hadrons use pseudorapidity), so they are kept separate. The new YAML
        # schema nests these under hadron:/jet: (hadron.eta, jet.eta, jet.eta_R, jet.rapidity);
        # the old top-level eta_cut/y_cut/eta_cut_R keys are kept as a fallback.
        _jet = block.get("jet") if isinstance(block.get("jet"), dict) else {}
        _hadron = block.get("hadron") if isinstance(block.get("hadron"), dict) else {}
        if "eta_cut" in block:
            self.eta_cut = block["eta_cut"]
        elif "eta" in _hadron:
            self.eta_cut = _hadron["eta"]
        elif "eta" in _jet:
            self.eta_cut = _jet["eta"]
        if "y_cut" in block:
            self.y_cut = block["y_cut"]
        elif "rapidity" in _jet:
            self.y_cut = _jet["rapidity"]
        # pt list may live at top-level (legacy), under jet:, or under hadron:
        if "pt" in block:
            self.pt = block["pt"]
        elif isinstance(block.get("jet"), dict) and "pt" in block["jet"]:
            self.pt = block["jet"]["pt"]
        elif isinstance(block.get("hadron"), dict) and "pt" in block["hadron"]:
            self.pt = block["hadron"]["pt"]
        # eta_R = jet outer pseudorapidity acceptance; the fiducial cut is eta_R - R. Runs after the
        # plain-eta block above so jets that declare eta_R get the R-subtracted value (matches the
        # legacy ordering where the eta_cut_R block came last).
        if "eta_cut_R" in block:
            self.eta_R = block["eta_cut_R"]
            self.eta_cut = np.round(self.eta_R - self.jet_R, decimals=1)
        elif "eta_R" in _jet:
            self.eta_R = _jet["eta_R"]
            self.eta_cut = np.round(self.eta_R - self.jet_R, decimals=1)
        if "c_ref" in block:
            index = block["jet"]["R"].index(self.jet_R)
            self.c_ref = block["c_ref"][index]
        if "low_trigger_range" in block:
            self.low_trigger_range = block["low_trigger_range"]
        if "high_trigger_range" in block:
            self.high_trigger_range = block["high_trigger_range"]
        if "trigger_range" in block:
            self.trigger_range = block["trigger_range"]
        self.logy = block.get("logy", False)
        # Whether the y-range / logy were set explicitly in the config. If not, the distribution
        # plotter auto-ranges the upper panel from the histogram+data content (see _auto_y_range),
        # since the new YAML schema rarely sets these and the old [0, 1.99] linear default left
        # many distributions invisible.
        self.logy_explicit = ("logy" in block) or ("logy_pp" in block)
        # Display-only: scale the pp MC so its integral matches the data over the common measured
        # range (shape comparison). For observables whose data is an absolute cross-section while the
        # MC is a per-event yield (e.g. pt_ch_atlas, mb GeV^-2), the two otherwise don't overlay.
        # R_AA stays absolute (this only rescales a display clone in plot_distribution_and_ratio).
        self.area_normalize = block.get("area_normalize", False)

        # for v2
        if "v2" in observable:
            self.y_ratio_min = -0.5
            self.y_ratio_max = 1.99

        if self.is_AA:
            # Default ytitle (the new schema comments out ytitle_AA on most blocks) so plot_RAA's
            # SetYTitle never sees an unset attribute; mirrors the pp branch's .get default.
            self.ytitle = block.get("ytitle_AA", "")
            self.y_range_explicit = "y_min_AA" in block
            if self.y_range_explicit:
                self.y_min = float(block["y_min_AA"])
                self.y_max = float(block["y_max_AA"])
            else:
                self.y_min = 0.0
                self.y_max = 1.0
            # Ratio-panel y-range. The AA branch never set these (only pp/v2 did), so plot_RAA's
            # SetRangeUser(self.y_ratio_min, self.y_ratio_max) hit an unset attribute. Same default
            # as pp; honour a per-block override if present.
            if "y_ratio_min" in block:
                self.y_ratio_min = block["y_ratio_min"]
                self.y_ratio_max = block["y_ratio_max"]
            else:
                self.y_ratio_min = 0.0
                self.y_ratio_max = 1.99
        else:
            self.ytitle = block.get("ytitle_pp", "")
            self.y_range_explicit = "y_min_pp" in block
            if self.y_range_explicit:
                self.y_min = float(block["y_min_pp"])
                self.y_max = float(block["y_max_pp"])
            else:
                self.y_min = 0.0
                self.y_max = 1.99
            if "y_ratio_min" in block:
                self.y_ratio_min = block["y_ratio_min"]
                self.y_ratio_max = block["y_ratio_max"]
            else:
                self.y_ratio_min = 0.0
                self.y_ratio_max = 1.99
            # Provide the opportunity to disable logy in pp, but keep in AA
            if "logy_pp" in block:
                self.logy = block["logy_pp"]

        self.skip_pp = block.get("skip_pp", False)
        self.skip_pp_ratio = block.get("skip_pp_ratio", False)
        self.skip_AA_ratio = block.get("skip_AA_ratio", False)
        # AA_difference: the measured AA modification is the *difference* of the per-jet
        # distributions (D_PbPb - D_pp), not the ratio -- e.g. CMS Dpt_cms's Delta_D(z),
        # whose data goes negative. When set, plot_RAA subtracts the pp reference instead of
        # dividing by it, keeps the config AA y-range (which must admit negatives), and draws
        # the reference line at 0 instead of 1.
        self.AA_difference = block.get("AA_difference", False)
        self.scale_by = block.get("scale_by", None)  # noqa: SIM910

        # Flag to plot hole histogram (for hadron histograms only)
        if self.is_AA:
            self.subtract_holes = observable in [
                "pt_ch_alice",
                "pt_pi_alice",
                "pt_pi0_alice",
                "pt_ch_cms",
                "pt_ch_atlas",
                "pt_pi0_phenix",
                "pt_ch_star",
                "v2_atlas",
                "v2_cms",
                "dihadron_star",
            ]
        else:
            self.subtract_holes = False

    # -------------------------------------------------------------------------------------------
    # Initialize JETSCAPE distribution into self.observable_settings
    # -------------------------------------------------------------------------------------------
    def initialize_jetscape_distribution(
        self, observable_type, observable, block, centrality, centrality_index, pt_suffix="", self_normalize=False
    ):
        # -------------------------------------------------------------
        # AA
        if self.is_AA:
            # Add centrality bin to list, if needed
            if centrality not in self.observable_centrality_list:
                self.observable_centrality_list.append(centrality)

            # -------------------------------------------------------------
            # For hadron observables, we retrieve and subtract the hole histograms
            if self.subtract_holes:
                hole_labels = ["", "_holes"]
                for hole_label in hole_labels:
                    # Get histogram, and add to self.observable_settings
                    self.get_histogram(
                        observable_type, observable, centrality, collection_label=hole_label, pt_suffix=pt_suffix
                    )

                    # Normalization
                    # Note: If we divide by the sum of weights (corresponding to n_events) and multiply by the
                    #       pt-hat cross-section, then JETSCAPE distribution gives cross-section: dsigma/dx (in mb)
                    self.scale_histogram(
                        observable_type,
                        observable,
                        centrality,
                        collection_label=hole_label,
                        pt_suffix=pt_suffix,
                        self_normalize=self_normalize,
                    )

                # Subtract the holes (and save unsubtracted histogram)
                if self.observable_settings["jetscape_distribution"]:
                    self.observable_settings["jetscape_distribution_unsubtracted"] = self.observable_settings[
                        "jetscape_distribution"
                    ].Clone()
                    self.observable_settings["jetscape_distribution_unsubtracted"].SetName(
                        "{}_unsubtracted".format(self.observable_settings["jetscape_distribution"].GetName())
                    )
                    if self.observable_settings["jetscape_distribution_holes"]:
                        self.observable_settings["jetscape_distribution"].Add(
                            self.observable_settings["jetscape_distribution_holes"], -1
                        )

                # Perform any additional manipulations on scaled histograms
                self.post_process_histogram(observable_type, observable, block, centrality, centrality_index)

            # -------------------------------------------------------------
            # For jet histograms, loop through all available hole subtraction variations, and initialize histogram
            else:
                for jet_collection_label in self.jet_collection_labels:
                    # Get histogram and add to self.observable_settings
                    #  - In the case of semi-inclusive measurements construct difference of histograms
                    self.get_histogram(
                        observable_type,
                        observable,
                        centrality,
                        pt_suffix=pt_suffix,
                        collection_label=jet_collection_label,
                    )
                    self.scale_histogram(
                        observable_type,
                        observable,
                        centrality,
                        collection_label=jet_collection_label,
                        pt_suffix=pt_suffix,
                        self_normalize=self_normalize,
                    )
                    self.post_process_histogram(
                        observable_type,
                        observable,
                        block,
                        centrality,
                        centrality_index,
                        collection_label=jet_collection_label,
                    )

        # -------------------------------------------------------------
        # pp
        else:
            self.get_histogram(observable_type, observable, centrality, pt_suffix=pt_suffix)
            self.scale_histogram(
                observable_type, observable, centrality, pt_suffix=pt_suffix, self_normalize=self_normalize
            )
            self.post_process_histogram(observable_type, observable, block, centrality, centrality_index)

    # -------------------------------------------------------------------------------------------
    # Build the encoder-based storage column name for a migrated jet/substructure observable.
    #
    # This mirrors, byte-for-byte, the name the analyzer writes (and the histogrammer histogrammed)
    # via obs.encode_name_for_storing_in_file(...). It is the COLUMN-name encoding: the grooming is
    # passed via convert_to_grooming_method_spec (the underlying SoftDropSpec/DynamicalGroomingSpec,
    # NO "SD_"/"DyG_" prefix) -- distinct from the DATA-BLOCK key encoding (GroomingSettingsSpec /
    # JetAxisDifferenceSpec .encode(), WITH the prefix) used only for data_block_params. Mirror of
    # HistogramResults._encoder_column_name; keep the two in lockstep.
    # -------------------------------------------------------------------------------------------
    def _encoder_column_name(
        self,
        observable_type,
        observable,
        jet_R,
        jet_collection_label,
        block=None,
        grooming_setting=None,
        pt_bin=None,
        axis_entry=None,
        charge=None,
        angularity=None,
    ):
        obs = self.observables_info[f"{self.sqrts}_{observable_type}_{observable}"]
        kwargs = {"jet_R": observable_module.JetRSpec(jet_R)}

        if pt_bin is not None and isinstance((block or {}).get("jet", {}).get("pt"), list):
            pt_edges = block["jet"]["pt"]
            if pt_bin + 1 < len(pt_edges):
                kwargs["jet_pt"] = observable_module.PtSpec(pt_edges[pt_bin], pt_edges[pt_bin + 1])

        if axis_entry is not None:
            # Grooming is encoded inside the jet-axis spec; do not also pass jet_grooming_settings.
            kwargs["jet_axis"] = observable_module.JetAxisDifferenceSpec(
                type=axis_entry["type"], grooming_settings=axis_entry.get("grooming_settings")
            )
        elif grooming_setting is not None:
            # Mirror the analyzer exactly: it passes the underlying method spec (no "SD_"/"DyG_"
            # prefix) as jet_grooming_settings for the column name.
            kwargs["jet_grooming_settings"] = observable_module.convert_to_grooming_method_spec(grooming_setting)

        # Step 4.5: jet charge (charge_cms) is parametrized by kappa, an essential parameter. The
        # analyzer encodes it as jet_charge=JetChargeSpec(kappa); the data block keys its tables with
        # the same JetChargeSpec encoding (kappa_X). Independent of grooming/axis (none coexist today).
        if charge is not None:
            kwargs["jet_charge"] = observable_module.JetChargeSpec(charge)

        # Step 5: angularity observables are parametrized by alpha (jet_angularity, essential). The
        # analyzer/histogrammer encode it as jet_angularity=AngularitySpec(alpha) (-> alpha_X_kappa_1.0);
        # `angularity` is an AngularitySpec already. Keep lockstep with HistogramResults._encoder_column_name.
        if angularity is not None:
            kwargs["jet_angularity"] = angularity

        return obs.encode_name_for_storing_in_file(tag=jet_collection_label, **kwargs)

    # -------------------------------------------------------------------------------------------
    # Get histogram and add to self.observable_settings
    #  - In AA case, also add hole histogram
    #  - In the case of semi-inclusive measurements construct difference of histograms
    # -------------------------------------------------------------------------------------------
    def get_histogram(self, observable_type, observable, centrality, collection_label="", pt_suffix=""):
        keys = [key.ReadObj().GetTitle() for key in self.input_file.GetListOfKeys()]

        # In the case of semi-inclusive recoil measurements, construct Delta_recoil from the
        # (aggregated) per-trigger-class yields. This MUST happen here, post-aggregation -- the
        # per-trigger normalization is non-linear (1/N_trig), so doing it per-file in the
        # histogrammer and then hadd-summing over-counts by n_files. observable_type is
        # "hadron_trigger_chjet" (IAA_pt_alice/dphi_alice) on 5020/2760; the legacy 2760/200
        # "semi_inclusive_*" naming is also routed here for back-compat.
        if "semi_inclusive" in observable_type or observable_type == "hadron_trigger_chjet":
            self.construct_semi_inclusive_histogram(
                keys, observable_type, observable, centrality, collection_label=collection_label, pt_suffix=pt_suffix
            )

        # For all other histograms, get the histogram directly
        else:
            # Get histogram. For the encoder-migrated jet/substructure observables, rebuild the
            # name from the same encoder the histogrammer used (recomputed per collection_label,
            # whose tag is folded into the encoder name) -- the pt range is already encoded, so the
            # _pt{i} suffix is suppressed here (it is still used by the binning/overlay lookups).
            if getattr(self, "_is_migrated_obs", False):
                column_name = self._encoder_column_name(
                    observable_type,
                    observable,
                    self.jet_R,
                    collection_label,
                    block=self._encoder_block,
                    grooming_setting=self._encoder_grooming_setting,
                    pt_bin=self._encoder_pt_bin,
                    axis_entry=self._encoder_axis_entry,
                    charge=self._encoder_charge,
                    angularity=self._encoder_angularity,
                )
                self.hname = f"h_{column_name}_{centrality}"
            else:
                self.hname = f"h_{observable_type}_{observable}{self.suffix}{collection_label}_{centrality}{pt_suffix}"
            if self.hname in keys:
                h_jetscape = self.input_file.Get(self.hname)
                h_jetscape.SetDirectory(0)
                if not h_jetscape.GetSumw2():
                    h_jetscape.Sumw2()

            else:
                h_jetscape = None
            self.observable_settings[f"jetscape_distribution{collection_label}"] = h_jetscape

            # pp run only: also fetch the R_AA-denominator histogram booked on the AA `ratio`
            # binning (same hname + "_raa_denom"). plot_RAA (AA run) divides by this to avoid the
            # spectra-vs-ratio bin-count mismatch. None when the histogrammer didn't book one (i.e.
            # ratio binning == spectra binning, or no ratio table) -- plot_RAA falls back then.
            h_raa_denom = None
            if not self.is_AA:
                raa_denom_name = f"{self.hname}_raa_denom"
                if raa_denom_name in keys:
                    h_raa_denom = self.input_file.Get(raa_denom_name)
                    h_raa_denom.SetDirectory(0)
                    if not h_raa_denom.GetSumw2():
                        h_raa_denom.Sumw2()
            self.observable_settings[f"jetscape_distribution_raa_denom{collection_label}"] = h_raa_denom

    # -------------------------------------------------------------------------------------------
    # Construct semi-inclusive observables from difference of histograms
    # -------------------------------------------------------------------------------------------
    def construct_semi_inclusive_histogram(self, keys, observable_type, observable, centrality, collection_label="", pt_suffix=""):
        # Semi-inclusive recoil (hadron_trigger_chjet: IAA_pt_alice / dphi_alice). Build Delta_recoil
        # from the AGGREGATED raw per-trigger-class yields and the AGGREGATED trigger-pt (N_trig)
        # counts: Delta_recoil = high/N_trig^high - c_ref * low/N_trig^low (each width-normalized).
        # The per-trigger normalization is done ONCE here, post-aggregation, so it is correct after
        # hadd-ing N files -- doing it per-file in the histogrammer then summing over-counts by n_files.
        if observable_type == "hadron_trigger_chjet":
            block = self.config[observable_type][observable]
            hname_high = f"h_{observable_type}_{observable}_R{self.jet_R}_highTrigger{collection_label}_{centrality}{pt_suffix}"
            hname_low = f"h_{observable_type}_{observable}_R{self.jet_R}_lowTrigger{collection_label}_{centrality}{pt_suffix}"
            # The analyzer stores the trigger-pt histogram once (under whichever observable owns it,
            # e.g. IAA_pt_alice); the trigger selection is shared, so find it by pattern.
            hname_ntrig = next(
                (k for k in keys if f"_trigger_pt{collection_label}" in k and k.endswith(f"_{centrality}")), None
            )
            self.hname = f"h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}{pt_suffix}"
            self.observable_settings[f"jetscape_distribution_raa_denom{collection_label}"] = None
            if hname_high in keys and hname_low in keys and hname_ntrig:
                h_high = self.input_file.Get(hname_high)
                h_high.SetDirectory(0)
                h_low = self.input_file.Get(hname_low)
                h_low.SetDirectory(0)
                h_ntrig = self.input_file.Get(hname_ntrig)
                h_ntrig.SetDirectory(0)
                trig = block.get("trigger", {})
                low_r = trig.get("low_range", block.get("low_trigger_range"))
                high_r = trig.get("high_range", block.get("high_trigger_range"))
                n_low = h_ntrig.GetBinContent(h_ntrig.FindBin(0.5 * (low_r[0] + low_r[1])))
                n_high = h_ntrig.GetBinContent(h_ntrig.FindBin(0.5 * (high_r[0] + high_r[1])))
                if n_low > 0 and n_high > 0:
                    c_ref = 1.0  # pp; per-R underlying-event correction for Pb-Pb
                    if self.is_AA and isinstance(block.get("c_ref"), list):
                        R_index = next(
                            (i for i, r in enumerate(block["jet"]["R"]) if np.isclose(r, self.jet_R)), 0
                        )
                        c_ref = block["c_ref"][R_index]
                    h_delta = h_high.Clone(self.hname)
                    h_delta.SetTitle(self.hname)
                    h_delta.SetDirectory(0)
                    h_delta.Scale(1.0 / n_high, "width")
                    h_low_n = h_low.Clone(f"{hname_low}_tmpnorm")
                    h_low_n.SetDirectory(0)
                    h_low_n.Scale(1.0 / n_low, "width")
                    h_delta.Add(h_low_n, -1.0 * c_ref)
                    # dphi only: the HEPData Delta_recoil(Dphi) is DOUBLE-differential -- per rad AND
                    # per recoil-jet-pt (GeV/c x rad)^-1. The MC histogram is binned in Dphi but
                    # INTEGRATED over the jet-pt window of this sub-bin, so divide by that window width
                    # to match the data (else it is high by the window, e.g. x10 for the 20-30 GeV bin;
                    # confirmed at full stats 2026-06-04). IAA's pt-spectrum data is Dphi-integrated, so
                    # it needs no analogous factor (it already matches at ratio ~1).
                    if "dphi" in observable and pt_suffix and isinstance(block.get("jet", {}).get("pt"), list):
                        pt_index = int(pt_suffix.replace("_pt", ""))
                        pt_edges = block["jet"]["pt"]
                        if pt_index + 1 < len(pt_edges):
                            pt_window = pt_edges[pt_index + 1] - pt_edges[pt_index]
                            if pt_window > 0:
                                h_delta.Scale(1.0 / pt_window)
                    self.observable_settings[f"jetscape_distribution{collection_label}"] = h_delta
                else:
                    self.observable_settings[f"jetscape_distribution{collection_label}"] = None
            else:
                self.observable_settings[f"jetscape_distribution{collection_label}"] = None
            return

        if self.sqrts == 2760:  # Delta recoil
            hname_low_trigger = (
                f"h_{observable_type}_{observable}_R{self.jet_R}_lowTrigger{collection_label}_{centrality}"
            )
            hname_high_trigger = (
                f"h_{observable_type}_{observable}_R{self.jet_R}_highTrigger{collection_label}_{centrality}"
            )
            hname_ntrigger = f"h_{observable_type}_alice_trigger_pt{observable}{collection_label}_{centrality}"
            self.hname = f"h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}"
            if hname_low_trigger in keys and hname_high_trigger in keys and hname_ntrigger in keys:
                h_jetscape_low = self.input_file.Get(hname_low_trigger)
                h_jetscape_low.SetDirectory(0)
                h_jetscape_high = self.input_file.Get(hname_high_trigger)
                h_jetscape_high.SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)

                low_trigger = (self.low_trigger_range[0] + self.low_trigger_range[1]) / 2
                high_trigger = (self.high_trigger_range[0] + self.high_trigger_range[1]) / 2
                n_trig_high = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(high_trigger))
                n_trig_low = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(low_trigger))
                h_jetscape_high.Scale(1.0 / n_trig_high)
                h_jetscape_low.Scale(1.0 / n_trig_low)
                self.observable_settings[f"jetscape_distribution{collection_label}"] = h_jetscape_high.Clone(
                    "h_delta_recoil"
                )
                self.observable_settings[f"jetscape_distribution{collection_label}"].Add(h_jetscape_low, -1)
            else:
                self.observable_settings[f"jetscape_distribution{collection_label}"] = None

        elif self.sqrts == 200:
            hname = f"h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}"
            hname_ntrigger = f"h_{observable_type}_star_trigger_pt{collection_label}_{centrality}"
            self.hname = f"h_{observable_type}_{observable}_R{self.jet_R}{collection_label}_{centrality}"
            if hname in keys and hname_ntrigger in keys:
                self.observable_settings[f"jetscape_distribution{collection_label}"] = self.input_file.Get(hname)
                self.observable_settings[f"jetscape_distribution{collection_label}"].SetDirectory(0)
                h_jetscape_ntrigger = self.input_file.Get(hname_ntrigger)
                h_jetscape_ntrigger.SetDirectory(0)

                trigger = (self.trigger_range[0] + self.trigger_range[1]) / 2
                n_trig = h_jetscape_ntrigger.GetBinContent(h_jetscape_ntrigger.FindBin(trigger))
                self.observable_settings[f"jetscape_distribution{collection_label}"].Scale(1.0 / n_trig)
            else:
                self.observable_settings[f"jetscape_distribution{collection_label}"] = None

    # -------------------------------------------------------------------------------------------
    # Scale histogram for:
    #  - xsec/weight_sum
    #  - bin width
    #  - observable-specific factors (eta, sigma_inel, n_jets, ...)
    #  - self-normalization
    # -------------------------------------------------------------------------------------------
    def scale_histogram(
        self,
        observable_type,
        observable: str,
        centrality,
        collection_label: str = "",
        pt_suffix: str = "",
        self_normalize: bool = False,
    ):
        # Apply the IDENTICAL scaling chain to the prediction histogram and, when present (pp run),
        # to the "_raa_denom" companion booked on the AA `ratio` binning. Scaling both through the
        # same path (xsec/weight_sum, bin width, observable-specific factors, self-normalization)
        # keeps the denominator a correctly-normalized pp spectrum so plot_RAA's Divide is valid.
        for settings_key in (
            f"jetscape_distribution{collection_label}",
            f"jetscape_distribution_raa_denom{collection_label}",
        ):
            if settings_key not in self.observable_settings:
                continue
            self._scale_one_histogram(
                self.observable_settings[settings_key],
                observable_type,
                observable,
                centrality,
                collection_label=collection_label,
                pt_suffix=pt_suffix,
                self_normalize=self_normalize,
            )

    # -------------------------------------------------------------------------------------------
    # Apply the full scaling chain to a single histogram (see scale_histogram).
    # -------------------------------------------------------------------------------------------
    def _scale_one_histogram(  # noqa: C901
        self,
        h,
        observable_type,
        observable: str,
        centrality,
        collection_label: str = "",
        pt_suffix: str = "",
        self_normalize: bool = False,
    ):
        if not h:
            return

        # (0) No scaling is needed for hadron v2 and jet v2
        if "v2" in observable:
            return

        # (0b) Semi-inclusive recoil (hadron_trigger_chjet) Delta_recoil observables (IAA_pt_alice,
        # dphi_alice) are already fully normalized in the histogrammer: the per-trigger ratio
        # (1/N_trig) cancels the xsec/weight_sum factor, and the bin-width scaling is already applied
        # (see _build_semi_inclusive_delta_recoil). Applying the standard chain again would
        # double-normalize, so skip it. But `not self_normalize` lets the self-normalized
        # hadron_trigger_chjet observable (nsubjettiness_alice, STAT_2760) fall through to the
        # self-normalization step below (it is not pre-normalized to a per-trigger scale).
        if observable_type == "hadron_trigger_chjet" and not self_normalize:
            return

        # --------------------------------------------------
        # (1) Scale all histograms by the min-pt-hat cross-section and weight-sum
        if self.is_AA:
            h_xsec = self.input_file.Get(f"h_xsec_{centrality}")
            h_weight_sum = self.input_file.Get(f"h_weight_sum_{centrality}")
            if not h_xsec or not h_weight_sum:
                logger.warning(f"Skipping {centrality} — missing h_xsec or h_weight_sum")
                return
            weight_sum = h_weight_sum.GetBinContent(1)
        else:
            h_xsec = self.input_file.Get("h_xsec")
            weight_sum = self.input_file.Get("h_weight_sum").GetBinContent(1)
        xsec = h_xsec.GetBinContent(1) / h_xsec.GetEntries()
        h.Scale(xsec / weight_sum)

        # --------------------------------------------------
        # (2) Scale all histograms by bin width
        h.Scale(1.0, "width")

        # --------------------------------------------------
        # (3) Scale for observable-specific factors
        if self.sqrts == 200:
            if observable_type == "hadron":
                if observable == "pt_ch_star":
                    sigma_inel = 42.0
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / (2 * np.pi))
                    h.Scale(1.0 / sigma_inel)
                if observable == "pt_pi0_phenix":
                    sigma_inel = 42.0  # Update
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / (2 * np.pi))
                    h.Scale(1.0 / sigma_inel)
            elif observable_type == "hadron_trigger_hadron":
                if observable == "dihadron_star":
                    # Need to grab the pt trig range
                    # Example: "_pt_trig_8_15_...." -> split -> ['', 'pt', 'trig', '8', '15', '....']
                    # So we grab index 3 and 4
                    split_suffix = self.suffix.split("_")
                    pt_trig_min, pt_trig_max = float(split_suffix[3]), float(split_suffix[4])

                    # And then the trigger hist
                    hname = f"h_{observable_type}_dihadron_star_Ntrig_{centrality}"
                    h_ntrig = self.input_file.Get(hname)
                    h_ntrig.SetDirectory(0)

                    # Finally, scale by the n_trig
                    n_trig = h_ntrig.GetBinContent(
                        h_ntrig.GetXaxis().FindBin((pt_trig_max + pt_trig_min) / 2)
                    )  # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_trig > 0.0:
                        h.Scale(1.0 / (n_trig * (xsec / weight_sum)))
                    else:
                        logger.warning("N_trig for dihadron_star == 0")
            elif observable_type == "inclusive_chjet":
                if observable == "pt_star":
                    h.Scale(1.0 / (2 * self.eta_cut))
            elif observable_type == "hadron_trigger_chjet":
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ["IAA_pt_star", "dphi_star"]:
                    h.Scale(1.0 / (xsec / weight_sum))

        if self.sqrts == 2760:
            if observable_type == "hadron":
                if observable in ["pt_ch_alice", "pt_pi_alice", "pt_pi0_alice"]:
                    sigma_inel = 61.8
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / sigma_inel)
                elif observable == "pt_ch_atlas":
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / (2 * np.pi))
                elif observable == "pt_ch_cms":
                    sigma_inel = 64.0
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / (2 * np.pi))
                    h.Scale(1.0 / sigma_inel)
            elif observable_type == "inclusive_jet":
                if observable == "pt_alice":
                    sigma_inel = 62.1
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / sigma_inel)
                if observable == "pt_atlas":
                    h.Scale(1.0 / (2 * self.y_cut))
                    h.Scale(1.0e6)  # convert to nb
                if observable == "pt_cms":
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0e6)  # convert to nb
                if observable in ["Dz_atlas", "Dpt_atlas", "Dz_cms", "Dpt_cms"]:
                    if observable in ["Dz_atlas", "Dpt_atlas"]:
                        hname = (
                            f"h_{observable_type}_Dz_atlas{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}"
                        )
                    elif observable in ["Dz_cms", "Dpt_cms"]:
                        hname = (
                            f"h_{observable_type}_Dz_cms{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}"
                        )
                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(
                        1
                    )  # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.0:
                        h.Scale(1.0 / (n_jets * (xsec / weight_sum)))
                    else:
                        logger.warning("N_jets = 0")
            elif observable_type == "inclusive_chjet":
                if observable == "pt_alice":
                    h.Scale(1.0 / (2 * self.eta_cut))
            elif observable_type == "hadron_trigger_chjet":  # noqa: SIM102
                # Note that n_trig histograms should also be scaled by xsec/n_events
                if observable in ["IAA_pt_alice", "dphi_alice"]:
                    h.Scale(1.0 / (xsec / weight_sum))

        if self.sqrts == 5020:
            if observable_type == "hadron":
                if observable in ["pt_ch_alice", "pt_pi_alice"]:
                    sigma_inel = 67.6
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / sigma_inel)
                elif observable == "pt_ch_cms":
                    sigma = 70.0
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0 / (2 * np.pi))
                    h.Scale(1.0 / sigma)
            elif observable_type == "inclusive_jet":
                if observable == "pt_alice":
                    h.Scale(1.0 / (2 * self.eta_cut))
                if observable in ["pt_atlas", "pt_y_atlas"]:
                    h.Scale(1.0e6)  # convert to nb
                    if observable == "pt_atlas":
                        h.Scale(1.0 / (2 * self.y_cut))
                if observable == "pt_cms":
                    h.Scale(1.0 / (2 * self.eta_cut))
                    h.Scale(1.0e6)  # convert to nb
                if observable in ["Dz_atlas", "Dpt_atlas"]:
                    hname = f"h_{observable_type}_Dz_atlas{self.suffix}{collection_label}_Njets_{centrality}{pt_suffix}"

                    h_njets = self.input_file.Get(hname)
                    h_njets.SetDirectory(0)
                    n_jets = h_njets.GetBinContent(
                        1
                    )  # Note that Njets histogram should also be scaled by xsec/n_events
                    if n_jets > 0.0:
                        h.Scale(1.0 / (n_jets * (xsec / weight_sum)))
                    else:
                        logger.warning("N_jets = 0")

        # Scale by bin-dependent factor (approximation for pp QA)
        if self.scale_by:
            nBins = h.GetNbinsX()
            for bin in range(1, nBins + 1):
                h_x = h.GetBinCenter(bin)
                h_y = h.GetBinContent(bin)
                h_error = h.GetBinError(bin)
                if self.scale_by == "1/pt":
                    scale_factor = 1.0 / h_x
                h.SetBinContent(bin, h_y * scale_factor)
                h.SetBinError(bin, h_error * scale_factor)

        # --------------------------------------------------
        # (4) Self-normalize histogram if appropriate
        if self_normalize:
            min_bin = 1
            if observable in ["zg_alice", "tg_alice"]:
                min_bin = 0
            nbins = h.GetNbinsX()
            integral = h.Integral(min_bin, nbins, "width")
            if integral > 0.0:
                h.Scale(1.0 / integral)
            else:
                logger.warning("integral for self-normalization is 0")

    # -------------------------------------------------------------------------------------------
    # Perform any additional manipulations on scaled histograms
    # -------------------------------------------------------------------------------------------
    def post_process_histogram(
        self, observable_type, observable, block, centrality, centrality_index: int, collection_label=""
    ):
        if "v2" in observable and self.is_AA:
            # hadron v2
            h = self.observable_settings[f"jetscape_distribution{collection_label}"]
            if h:
                h_num_name = f"h_{observable_type}_{observable}_{centrality}"
                h_num_name_holes = f"h_{observable_type}_{observable}_holes_{centrality}"
                h_denom_name = f"h_{observable_type}_{observable}_denom_{centrality}"
                h_denom_name_holes = f"h_{observable_type}_{observable}_holes_denom_{centrality}"

                self.observable_settings["jetscape_distribution"] = self.input_file.Get(h_num_name).Clone()
                self.observable_settings["jetscape_distribution_unsubtracted"] = self.input_file.Get(h_num_name).Clone()

                for i in range(0, self.input_file.Get(h_num_name).GetNbinsX()):  # noqa: PIE808
                    h_num_i = self.input_file.Get(h_num_name).GetBinContent(i)
                    h_num_holes_i = self.input_file.Get(h_num_name_holes).GetBinContent(i)
                    h_denom_i = self.input_file.Get(h_denom_name).GetBinContent(i)
                    h_denom_holes_i = self.input_file.Get(h_denom_name_holes).GetBinContent(i)
                    if h_denom_i == 0.0:
                        h_denom_i = -99.0
                        h_denom_holes_i = 0.0
                    self.observable_settings["jetscape_distribution"].SetBinContent(
                        i, (h_num_i - h_num_holes_i) / (h_denom_i - h_denom_holes_i)
                    )
                    self.observable_settings["jetscape_distribution_unsubtracted"].SetBinContent(i, h_num_i / h_denom_i)

                self.observable_settings["jetscape_distribution"].Print()
                self.observable_settings["jetscape_distribution_unsubtracted"].Print()

                self.observable_settings["jetscape_distribution"].SetName(
                    f"jetscape_distribution_{observable_type}_{observable}_{centrality}"
                )
                self.observable_settings["jetscape_distribution_unsubtracted"].SetName(
                    f"jetscape_distribution_unsubtracted_{observable_type}_{observable}_{centrality}"
                )

        # For the ATLAS rapidity-dependence, we need to divide the histograms by their first bin (|y|<0.3) to form a double ratio
        if observable == "pt_y_atlas":
            # Get y=0.3 entry, and create a histogram to divide by (so that uncertainties are propagated)
            h = self.observable_settings[f"jetscape_distribution{collection_label}"]
            if h:
                denominator = h.GetBinContent(1)
                denominator_uncertainty = h.GetBinError(1)
                h_denominator = h.Clone()
                h_denominator.SetName(f"{h_denominator.GetName()}_denominator")
                n = h_denominator.GetNbinsX()
                for i in range(n):
                    h_denominator.SetBinContent(i + 1, denominator)
                    h_denominator.SetBinError(i + 1, denominator_uncertainty)

                # Divide histograms
                h.Divide(h_denominator)

                # Rebin to remove first bin (|y|<0.3)
                bins = np.array(h.GetXaxis().GetXbins())
                hname = h.GetName()
                self.observable_settings[f"jetscape_distribution{collection_label}"] = h.Rebin(n - 1, hname, bins[1:])

        # In the case of pp, just look at the deltaPhi correlations themselves
        # In any case, we don't have data to compare against.
        # If skip_pp is turned off, then we can check the yields in pp.
        if observable_type == "hadron_trigger_hadron" and observable == "dihadron_star" and not self.skip_pp:
            # Grab the delta phi correlation
            h = self.observable_settings["jetscape_distribution"]
            # Extract the yield
            yield_lower_range_value, yield_upper_range_value = block["yield_range"]
            # First, near side
            _temp_error = ctypes.c_double(0)
            near_side_yield = h.IntegralAndError(
                h.GetXaxis().FindBin(yield_lower_range_value),
                h.GetXaxis().FindBin(yield_upper_range_value),
                _temp_error,
            )
            near_side_yield_error = _temp_error.value
            # Then the away side
            _temp_error = ctypes.c_double(0)
            away_side_yield = h.IntegralAndError(
                # Both are np.pi + because the yield_lower_range_value is negative
                h.GetXaxis().FindBin(np.pi + yield_lower_range_value),
                h.GetXaxis().FindBin(np.pi + yield_upper_range_value),
                _temp_error,
            )
            away_side_yield_error = _temp_error.value

            # Encode the values into a single TH1F. We put it into a TH1F instead of a TGraph because
            # the framework seems to require that the jetscape data is in a histogram.
            yield_bins = np.array(block["yield_bins"])
            n = len(yield_bins) - 1
            h_yield = ROOT.TH1F(f"{h.GetName()}_yield", f"{h.GetName()}_yields", n, yield_bins)

            # Fill in extracted values
            yield_values = np.zeros(n)
            yield_errors = np.zeros(n)
            yield_values[centrality_index] = near_side_yield
            yield_errors[centrality_index] = near_side_yield_error
            # round is just to ensure we have an int
            yield_values[centrality_index + round(n / 2)] = away_side_yield
            yield_errors[centrality_index + round(n / 2)] = away_side_yield_error

            # start at 1 to account for ROOT bin indexing
            for i, (yv, ye) in enumerate(zip(yield_values, yield_errors, strict=True), start=1):
                h_yield.SetBinContent(i, yv)
                h_yield.SetBinError(i, ye)

            # Now, put it together and store the result
            self.observable_settings["jetscape_distribution"] = h_yield

    # -------------------------------------------------------------------------------------------
    # Plot distributions
    # -------------------------------------------------------------------------------------------
    def plot_observable(self, observable_type, observable: str, centrality, pt_suffix: str = "") -> None:
        label = f"{observable_type}_{observable}{self.suffix}_{centrality}{pt_suffix}"

        if "v2" in observable:
            # for hadron v2
            if self.observable_settings["jetscape_distribution"]:
                self.plot_distribution_and_ratio(observable_type, observable, centrality, label, pt_suffix=pt_suffix)
            return

        # If AA: Plot PbPb/pp ratio, and comparison to data
        # If pp: Plot distribution, and ratio to data
        if self.is_AA:
            self.plot_RAA(observable_type, observable, centrality, label, pt_suffix=pt_suffix)
            self.write_experimental_data(label)
        elif self.observable_settings["jetscape_distribution"]:
            self.plot_distribution_and_ratio(observable_type, observable, centrality, label, pt_suffix=pt_suffix)

    # -------------------------------------------------------------------------------------------
    # Write experimental data tables
    # -------------------------------------------------------------------------------------------
    def write_experimental_data(self, label: str) -> None:
        output_dir = (self.output_dir / "../Data").resolve()
        output_dir.mkdir(exist_ok=True)

        force_write = True
        filename = f"Data_{label}.dat"
        output_file = output_dir / filename
        if force_write or not output_file.exists():
            # Get histogram binning ("raa_denom" is the R_AA-denominator helper hist, not a prediction)
            key = next(
                (
                    key
                    for key in self.observable_settings
                    if "jetscape_distribution" in key and "holes" not in key and "raa_denom" not in key
                )
            )
            h_prediction = self.observable_settings[key]

            # Truncate data tgraph to prediction histogram binning
            if h_prediction:
                xbins = np.array(h_prediction.GetXaxis().GetXbins())
                x_min = xbins[:-1]
                x_max = xbins[1:]

                g_data = self.observable_settings["data_distribution"]
                if not g_data:
                    return
                g_truncated = self.plot_utils.truncate_tgraph(g_data, h_prediction, is_AA=self.is_AA)
                if g_truncated:
                    y = np.array(g_truncated.GetY())
                    y_err = np.array([g_truncated.GetErrorY(i) for i in range(g_truncated.GetN())])
                    x = np.array(g_truncated.GetX())

                    # The Data_*.dat table pairs the prediction-histogram bins (x_min/x_max) with the
                    # truncated data points (x/y). If they don't line up -- different counts, or a point
                    # outside its bin -- the measured data simply doesn't cover this histogram's binning
                    # (a known HEPData edge case; see OBSERVABLE_EDGE_CASES). The R_AA plot is already
                    # rendered by now, so skip this ancillary curation table instead of aborting the run.
                    if len(x) != len(x_min) or np.any(np.greater(x_min, x)) or np.any(np.greater(x, x_max)):
                        logger.warning(
                            f"Skipping Data table {filename}: data points don't align with hist binning "
                            f"(x={x} vs. x_bins={xbins})"
                        )
                        return

                    df = pd.DataFrame({"x_min": x_min, "x_max": x_max, "y": y, "y_err": y_err})
                    # Write table
                    header = "Version 1.1\n"
                    header += "Label xmin xmax y y_err"
                    np.savetxt(output_file, df.values, header=header)
                else:
                    logger.info(f"Did not write Data table: {filename} -- missing tgraph")
            else:
                logger.info(f"Did not write Data table: {filename} -- missing histogram")
        else:
            logger.info(f"No need to write Data table: {filename} -- already exists")

    # -------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    # We plot all available modes of hole subtraction:
    #   For hadron histograms (self.subtract_holes = True):
    #       - self.observable_settings['jetscape_distribution']
    #       - self.observable_settings['jetscape_distribution_unsubtracted']
    #   For jet histograms:
    #       for jet_collection_label in self.jet_collection_labels:
    #           self.observable_settings[f'jetscape_distribution{jet_collection_label}']
    # -------------------------------------------------------------------------------------------
    def plot_RAA(self, observable_type, observable, centrality, label, pt_suffix=""):  # noqa: C901
        # Assemble the list of hole subtraction variations ("raa_denom" is the pp-only
        # R_AA-denominator helper hist, never a prediction to plot/divide here)
        keys_to_plot = [
            key
            for key in self.observable_settings
            if "jetscape_distribution" in key and "holes" not in key and "raa_denom" not in key
            # For hadron observables (subtract_holes), plot only the hole-subtracted (physical)
            # R_AA, not the unsubtracted cross-check: the two nearly coincide at the measured pT,
            # so drawing both leaves a single visible band whose 2nd-entry color clashes with the
            # primary "JETSCAPE" legend entry. Jets keep all hole-subtraction variations.
            and not (self.subtract_holes and "unsubtracted" in key)
        ]
        model_display_name = _model_display_name[self.model_name]
        self.jetscape_legend_label = {}
        self.jetscape_legend_label["jetscape_distribution"] = model_display_name
        self.jetscape_legend_label["jetscape_distribution_unsubtracted"] = f"{model_display_name} (unsubtracted)"
        self.jetscape_legend_label["jetscape_distribution_shower_recoil"] = f"{model_display_name} (shower+recoil)"
        self.jetscape_legend_label["jetscape_distribution_shower_recoil_unsubtracted"] = (
            f"{model_display_name} (shower+recoil, unsubtracted)"
        )
        self.jetscape_legend_label["jetscape_distribution_negative_recombiner"] = (
            f"{model_display_name} (negative recombiner)"
        )
        self.jetscape_legend_label["jetscape_distribution_constituent_subtraction"] = f"{model_display_name} (CS)"

        if not self.observable_settings[keys_to_plot[0]]:
            logger.warning(f"Skipping {label} since data is missing")
            return

        # A self-normalized DISTRIBUTION comparison (skip_AA_ratio) exists only to overlay the
        # measured PbPb distribution against the MC. When this (jet_R, pt, ...) bin has no published
        # HEPData distribution -- e.g. axis_alice R=0.4 outside pt 80-100, where ALICE measured no
        # PbPb distribution -- there is nothing to compare, so skip the plot rather than emit a
        # data-less MC-only panel. (Genuine R_AA plots, skip_AA_ratio=False, still draw without an
        # overlay: the MC R_AA is meaningful on its own.)
        if self.skip_AA_ratio and not self.observable_settings.get("data_distribution"):
            logger.info(f"Skipping {label}: skip_AA_ratio distribution with no measured PbPb data for this bin")
            return

        # Get the pp reference histogram and form the RAA ratios.
        # Prefer the "_raa_denom" pp histogram (booked on the AA `ratio` binning) so the bin edges
        # match the AA arm -- this avoids the "Cannot divide histograms with different number of
        # bins" error when the pp spectra binning differs from the ratio binning. Fall back to the
        # plain spectra histogram when no _raa_denom was produced (ratio binning == spectra binning,
        # or no ratio table), which preserves the previous behavior for those observables.
        if not self.skip_AA_ratio:
            h_pp_name = f"jetscape_distribution_raa_denom_{label}"
            h_pp = self.pp_ref_file.Get(h_pp_name)
            if not h_pp:
                h_pp_name = f"jetscape_distribution_{label}"
                h_pp = self.pp_ref_file.Get(h_pp_name)
            for key in keys_to_plot:
                if self.observable_settings[key] and h_pp:
                    if self.AA_difference:
                        # Measured quantity is D_PbPb - D_pp (a difference), not a ratio.
                        self.observable_settings[key].Add(h_pp, -1)
                    else:
                        self.observable_settings[key].Divide(h_pp)

        # Content extrema (data + MC, incl. errors) -- drives non-ratio auto-scaling AND the log twin.
        _objs = [self.observable_settings.get("data_distribution")] + [
            self.observable_settings[k] for k in keys_to_plot
        ]
        vmin, vposmin, vmax = self._content_extrema([o for o in _objs if o])

        # Y-range. True R_AA ratios keep the fixed user-requested [0, 3] (legend clears the points,
        # unity line visible). Non-ratio plots -- skip_AA_ratio distribution overlays and AA_difference
        # (D_PbPb - D_pp, which goes negative) -- auto-scale to the actual drawn content with extra
        # top headroom so the curves clear the legend/title rather than being squished into [0, 3].
        if self.skip_AA_ratio or self.AA_difference:
            y_lo, y_hi, _ = self._auto_y_range(vmin, vposmin, vmax, prefer_log=False)
            # _auto_y_range gives ~1.4x top headroom; bump to ~2x so the upper-panel legend clears.
            self.y_min, self.y_max = y_lo, y_lo + (y_hi - y_lo) * 1.45
        else:
            self.y_min, self.y_max = 0.0, 3.0

        c = ROOT.TCanvas("c", "c", 600, 450)
        c.SetRightMargin(0.05)
        c.SetLeftMargin(0.15)
        c.SetTopMargin(0.05)
        c.SetBottomMargin(0.17)
        c.cd()

        legend = ROOT.TLegend(0.4, 0.55, 0.75, 0.81)
        self.plot_utils.setup_legend(legend, 0.045, sep=-0.1)

        if not self.skip_AA_ratio:
            if not h_pp:
                logger.warning(f"Histogram {h_pp_name} not found or not valid. Skipping.")
                return
            self.bins = np.array(h_pp.GetXaxis().GetXbins())
        else:
            self.bins = np.array(self.observable_settings[keys_to_plot[0]].GetXaxis().GetXbins())

        myBlankHisto = ROOT.TH1F("myBlankHisto", "Blank Histogram", 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min)
        myBlankHisto.GetXaxis().SetTitleSize(0.07)
        myBlankHisto.GetYaxis().SetTitleSize(0.07)
        myBlankHisto.GetXaxis().SetTitleOffset(1.0)
        myBlankHisto.GetYaxis().SetTitleOffset(1.0)
        myBlankHisto.Draw("E")

        # Draw data
        if self.observable_settings["data_distribution"]:
            self.output_dict[f"data_distribution_{label}"] = self.observable_settings["data_distribution"]
            self.observable_settings["data_distribution"].SetName(f"data_distribution_{label}")
            self.observable_settings["data_distribution"].SetMarkerSize(self.marker_size)
            self.observable_settings["data_distribution"].SetMarkerStyle(self.data_marker)
            self.observable_settings["data_distribution"].SetMarkerColor(self.data_color)
            self.observable_settings["data_distribution"].SetLineStyle(self.line_style)
            self.observable_settings["data_distribution"].SetLineWidth(self.line_width)
            self.observable_settings["data_distribution"].SetLineColor(self.data_color)
            self.observable_settings["data_distribution"].Draw("PE Z same")
            legend.AddEntry(self.observable_settings["data_distribution"], "Data", "PE")

        # Draw JETSCAPE
        for i, key in enumerate(keys_to_plot):
            self.output_dict[f"{key}_{label}"] = self.observable_settings[key]
            if self.observable_settings[key]:
                if self.observable_settings[key].GetNbinsX() > 1:
                    self.observable_settings[key].SetFillColor(self.jetscape_color[i])
                    self.observable_settings[key].SetFillColorAlpha(self.jetscape_color[i], self.jetscape_alpha[i])
                    self.observable_settings[key].SetFillStyle(self.jetscape_fillstyle[i])
                    self.observable_settings[key].SetMarkerSize(0.0)
                    self.observable_settings[key].SetMarkerStyle(0)
                    self.observable_settings[key].SetLineWidth(0)
                    self.observable_settings[key].DrawCopy("E3 same")
                elif self.observable_settings[key].GetNbinsX() == 1:
                    self.observable_settings[key].SetMarkerSize(self.marker_size)
                    self.observable_settings[key].SetMarkerStyle(self.data_marker + 1)
                    self.observable_settings[key].SetMarkerColor(self.jetscape_color[i])
                    self.observable_settings[key].SetLineStyle(self.line_style)
                    self.observable_settings[key].SetLineWidth(self.line_width)
                    self.observable_settings[key].SetLineColor(self.jetscape_color[i])
                    self.observable_settings[key].DrawCopy("PE same")
                legend.AddEntry(self.observable_settings[key], self.jetscape_legend_label[key], "f")

        legend.Draw()

        # Reference line: zero for difference observables (D_PbPb - D_pp), unity for ratios.
        if not self.skip_AA_ratio and (self.AA_difference or self.y_max > 1):
            ref_y = 0.0 if self.AA_difference else 1.0
            line = ROOT.TLine(self.bins[0], ref_y, self.bins[-1], ref_y)
            line.SetLineColor(920 + 2)
            line.SetLineStyle(2)
            line.SetLineWidth(2)
            line.Draw()

        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.18
        text_latex.SetTextSize(0.065)
        text = f"#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts / 1000.0} TeV"
        text_latex.DrawLatex(x, 0.88, text)
        text = f"{centrality} {self.suffix} {pt_suffix}"
        text_latex.DrawLatex(x, 0.8, text)

        # For migrated observables use the unique encoder-based name (self.hname, set in
        # get_histogram) so per-grooming / per-axis variants don't collide on the legacy suffix
        # (e.g. axis_alice's two WTA_SD variants share `_R{R}_WTA_SD`). Legacy path unchanged.
        if getattr(self, "_is_migrated_obs", False) and self.hname:
            hname = self.hname
        else:
            hname = f"h_{observable_type}_{observable}{self.suffix}_{centrality}{pt_suffix}"
        c.SaveAs(str(self.output_dir / f"{hname}{self.file_format}"))
        # Log-y twin ("<name>_logxy.pdf") for non-negative AA plots -- skip AA_difference plots, which
        # go negative (log undefined). Reuses the distribution log-twin helper (log-x too if x>0).
        if not self.AA_difference and vposmin and vposmin > 0:
            self._save_logxy_twin(
                c, hname, dist_pad=c, blank_histo=myBlankHisto,
                log_ymin=vposmin / 3.0, log_ymax=vmax * 8.0,
            )
        c.Close()

    # -------------------------------------------------------------------------------------------
    # Save a log-x/log-y twin of a (non-ratio) distribution canvas.
    #
    # For each distribution/spectrum plot we additionally emit a "<name>_logxy.pdf" with log x (and
    # log y) on the upper distribution pad and log x on the ratio pad. This is purely a second view;
    # the pad log-flags are restored afterwards so the normal save (if it happens after) is unchanged.
    # NOT used for R_AA ratios (y ~ 0-2 with y-min often 0 -> degenerate log-y).
    # -------------------------------------------------------------------------------------------
    def _save_logxy_twin(
        self, canvas, base_name, dist_pad=None, ratio_pad=None, blank_histo=None, log_ymin=None, log_ymax=None
    ):
        # Remember current log state so we can restore it (callers may reuse the canvas/pads).
        saved = [(p, p.GetLogx(), p.GetLogy()) for p in (dist_pad, ratio_pad) if p is not None]
        # The main plot may be on a LINEAR y-range (so its y_min can be 0, which is invalid for
        # log-y). When a positive log-range is supplied, swap the blank histo onto it for the twin
        # so the curve isn't clipped/degenerate; restore afterwards.
        saved_range = None
        if blank_histo is not None and log_ymin and log_ymax and log_ymin > 0:
            saved_range = (blank_histo.GetMinimum(), blank_histo.GetMaximum())
            blank_histo.SetMinimum(log_ymin)
            blank_histo.SetMaximum(log_ymax)
        # Only use log-x where the x-range is strictly positive; observables whose x starts at 0
        # (e.g. groomed mass m_g, z_g) would render an empty log-x plot, so fall back to log-y only.
        can_logx = blank_histo is not None and blank_histo.GetXaxis().GetXmin() > 0
        try:
            if dist_pad is not None:
                if can_logx:
                    dist_pad.SetLogx()
                dist_pad.SetLogy()
            if ratio_pad is not None and can_logx:
                ratio_pad.SetLogx()
            canvas.Update()
            canvas.SaveAs(str(self.output_dir / f"{base_name}_logxy{self.file_format}"))
        finally:
            for p, logx, logy in saved:
                p.SetLogx(logx)
                p.SetLogy(logy)
            if saved_range is not None:
                blank_histo.SetMinimum(saved_range[0])
                blank_histo.SetMaximum(saved_range[1])
            canvas.Update()

    # -------------------------------------------------------------------------------------------
    # Content-driven y-axis range for the upper (distribution) panel.
    #
    # The new YAML schema rarely sets y_min/y_max/logy, so the old per-arm defaults ([0, 1.99]
    # linear) left many distributions invisible (curve pinned at the axis). Instead we derive the
    # range from the actual drawn content (JETSCAPE histograms + data graph, including errors) and
    # auto-pick log vs linear from the dynamic range. Explicit config ranges still win.
    # -------------------------------------------------------------------------------------------
    def _content_extrema(self, objs):
        """(vmin, vpos_min, vmax) over the drawn content of TH1/TGraph objects incl. errors; Nones if empty."""
        vmin = vposmin = vmax = None
        for o in objs:
            if not o:
                continue
            if isinstance(o, ROOT.TH1):
                for i in range(1, o.GetNbinsX() + 1):
                    c = o.GetBinContent(i)
                    e = o.GetBinError(i)
                    if c == 0 and e == 0:
                        continue
                    hi, lo = c + e, c - e
                    vmax = hi if vmax is None else max(vmax, hi)
                    vmin = lo if vmin is None else min(vmin, lo)
                    if c > 0:
                        vposmin = c if vposmin is None else min(vposmin, c)
            else:  # TGraph / TGraphErrors / TGraphAsymmErrors
                y = o.GetY()
                for i in range(o.GetN()):
                    yi = y[i]
                    ehi = max(o.GetErrorYhigh(i), 0.0)  # -1 for a plain TGraph -> 0
                    elo = max(o.GetErrorYlow(i), 0.0)
                    hi, lo = yi + ehi, yi - elo
                    vmax = hi if vmax is None else max(vmax, hi)
                    vmin = lo if vmin is None else min(vmin, lo)
                    if yi > 0:
                        vposmin = yi if vposmin is None else min(vposmin, yi)
        return vmin, vposmin, vmax

    def _log_yrange(self, vposmin, vmax):
        """Log-scale (y_min, y_max) with span-adaptive top headroom. Steeply-falling pT spectra span
        many decades and need far more headroom to clear the title/legend than O(1)-shape
        distributions, so scale the top factor by the dynamic range; ~half a decade at the bottom."""
        span = (vmax / vposmin) if (vposmin and vposmin > 0) else 1.0
        top = 1.0e4 if span >= 1.0e4 else 50.0  # pT spectra (>=4 decades) -> x1e4; others -> x50
        return vposmin / 3.0, vmax * top

    def _auto_y_range(self, vmin, vposmin, vmax, prefer_log=None):
        """Return (y_min, y_max, use_log) for the given content extrema (with headroom for the legend)."""
        if vmax is None or vmax <= 0:
            return 0.0, 1.0, bool(prefer_log)
        if vposmin is None:
            vposmin = vmax
        span = (vmax / vposmin) if vposmin > 0 else 1.0
        # >~2 decades of dynamic range -> log (unless the config forces a choice via prefer_log).
        use_log = prefer_log if prefer_log is not None else (span >= 100.0 and (vmin is None or vmin >= 0))
        if use_log:
            y_lo, y_hi = self._log_yrange(vposmin, vmax)
            return y_lo, y_hi, True
        y_lo = 0.0 if (vmin is None or vmin >= 0) else vmin - 0.1 * abs(vmax)
        return y_lo, vmax * 1.4, False

    def _area_norm_scale(self, h, g):
        """Factor to scale MC hist h so its integral matches data graph g over the bins both
        populate (sum of value x bin-width over matched bins). None if not computable."""
        mc_int = data_int = 0.0
        for i in range(g.GetN()):
            gx, gy, _yl, _yu = self.plot_utils.get_gx_gy(g, i)
            if gy == 0.0:
                continue
            b = h.FindBin(gx)
            if b < 1 or b > h.GetNbinsX():
                continue
            if not (h.GetXaxis().GetBinLowEdge(b) <= gx <= h.GetXaxis().GetBinUpEdge(b)):
                continue
            mc_y = h.GetBinContent(b)
            if mc_y == 0.0:
                continue  # below the MC kinematic cut (e.g. pT < 5 GeV) -> exclude; integrate only
                # over the common range where BOTH MC and data have content.
            w = h.GetXaxis().GetBinWidth(b)
            mc_int += mc_y * w
            data_int += gy * w
        if mc_int > 0.0 and data_int > 0.0:
            return data_int / mc_int
        return None

    # -------------------------------------------------------------------------------------------
    # Plot distributions in upper panel, and ratio in lower panel
    # -------------------------------------------------------------------------------------------
    def plot_distribution_and_ratio(  # noqa: C901
        self, observable_type, observable: str, centrality, label, pt_suffix: str = ""
    ):
        # Setup
        model_display_name = _model_display_name[self.model_name]

        c = ROOT.TCanvas("c", "c", 600, 650)
        c.Draw()
        c.cd()

        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad("myPad", "The pad", 0, pad2_dy, 1, 1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.0)
        pad1.SetTicks(0, 1)
        pad1.Draw()
        if self.logy:
            pad1.SetLogy()
        pad1.cd()

        legend = ROOT.TLegend(0.58, 0.6, 0.75, 0.75)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1)

        legend_ratio = ROOT.TLegend(0.25, 0.8, 0.45, 1.07)
        self.plot_utils.setup_legend(legend_ratio, 0.07, sep=-0.1)

        self.bins = np.array(self.observable_settings["jetscape_distribution"].GetXaxis().GetXbins())

        # Display-only area normalization (e.g. pt_ch_atlas): rescale a CLONE of the pp MC so its
        # integral matches the data over the common measured range, for a shape comparison. The
        # ORIGINAL is persisted to final_results.root (so R_AA stays absolute); only the drawn
        # histogram and the MC/data ratio in this plot use the scaled clone.
        jetscape_original = self.observable_settings["jetscape_distribution"]
        self._area_normalized = False
        if getattr(self, "area_normalize", False) and self.observable_settings.get("data_distribution"):
            scale = self._area_norm_scale(jetscape_original, self.observable_settings["data_distribution"])
            if scale and scale > 0:
                disp = jetscape_original.Clone(f"{jetscape_original.GetName()}_areanorm")
                disp.SetDirectory(0)
                disp.Scale(scale)
                self.observable_settings["jetscape_distribution"] = disp
                self.observable_settings["ratio"] = self.plot_utils.divide_histogram_by_tgraph(
                    disp, self.observable_settings["data_distribution"], include_tgraph_uncertainties=False
                )
                self._area_normalized = True
            else:
                logger.warning(
                    f"area_normalize requested for {observable_type}_{observable} but no overlapping "
                    f"MC/data range (scale={scale}) -- drawing un-normalized."
                )
        elif getattr(self, "area_normalize", False):
            logger.warning(
                f"area_normalize requested for {observable_type}_{observable} but no data graph -- not applied."
            )

        # Auto-range the upper panel from the drawn content (unless the config set y_min/y_max
        # explicitly). Also compute a guaranteed-positive log range for the _logxy twin.
        dist_objs = [self.observable_settings.get("data_distribution")]
        dist_objs += [
            v
            for k, v in self.observable_settings.items()
            if k.startswith("jetscape_distribution") and "raa_denom" not in k and "holes" not in k
        ]
        dist_objs = [o for o in dist_objs if o]
        _vmin, _vposmin, _vmax = self._content_extrema(dist_objs)
        if not getattr(self, "y_range_explicit", False) and _vmax and _vmax > 0:
            self.y_min, self.y_max, _auto_log = self._auto_y_range(
                _vmin, _vposmin, _vmax, prefer_log=(self.logy if getattr(self, "logy_explicit", False) else None)
            )
            if not getattr(self, "logy_explicit", False) and _auto_log:
                self.logy = True
                pad1.SetLogy()
        if _vposmin and _vmax and _vposmin > 0:
            log_ymin, log_ymax = self._log_yrange(_vposmin, _vmax)
        else:
            log_ymin = log_ymax = None

        myBlankHisto = ROOT.TH1F("myBlankHisto", "Blank Histogram", 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle(self.xtitle)
        myBlankHisto.SetYTitle(self.ytitle)
        myBlankHisto.SetMaximum(self.y_max)
        myBlankHisto.SetMinimum(self.y_min)  # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw("E")

        # Ratio
        c.cd()
        pad2 = ROOT.TPad("pad2", "pad2", 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0, 1)
        pad2.Draw()
        pad2.cd()

        myBlankHisto2 = myBlankHisto.Clone("myBlankHisto_C")
        myBlankHisto2.SetYTitle("Ratio")
        myBlankHisto2.SetXTitle(self.xtitle)
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.0)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(self.y_ratio_min, self.y_ratio_max)
        myBlankHisto2.Draw("")

        # Draw distribution
        pad1.cd()
        if self.observable_settings["data_distribution"]:
            self.output_dict[f"data_distribution_{label}"] = self.observable_settings["data_distribution"]
            self.observable_settings["data_distribution"].SetName(f"data_distribution_{label}")
            self.observable_settings["data_distribution"].SetMarkerSize(self.marker_size)
            self.observable_settings["data_distribution"].SetMarkerStyle(self.data_marker)
            self.observable_settings["data_distribution"].SetMarkerColor(self.data_color)
            self.observable_settings["data_distribution"].SetLineStyle(self.line_style)
            self.observable_settings["data_distribution"].SetLineWidth(self.line_width)
            self.observable_settings["data_distribution"].SetLineColor(self.data_color)
            self.observable_settings["data_distribution"].Draw("PE Z same")
            legend.AddEntry(self.observable_settings["data_distribution"], "Data", "PE")

        # Draw JETSCAPE
        # Persist the ORIGINAL (unscaled) MC -- even when area-normalized for display -- so the AA
        # R_AA division (which may fall back to this hist) stays on the absolute normalization.
        self.output_dict[f"jetscape_distribution_{label}"] = jetscape_original
        # Persist the R_AA-denominator companion (pp run only) into final_results.root so the AA
        # plot_RAA run can divide by it (it is scaled/normalized identically; only the binning is the
        # AA `ratio` binning). Not drawn here -- it's an internal histogram for the ratio, not a pp curve.
        h_raa_denom = self.observable_settings.get("jetscape_distribution_raa_denom")
        if h_raa_denom:
            self.output_dict[f"jetscape_distribution_raa_denom_{label}"] = h_raa_denom
        if self.observable_settings["jetscape_distribution"].GetNbinsX() > 1:
            self.observable_settings["jetscape_distribution"].SetFillColor(self.jetscape_color[0])
            self.observable_settings["jetscape_distribution"].SetFillColorAlpha(
                self.jetscape_color[0], self.jetscape_alpha[0]
            )
            self.observable_settings["jetscape_distribution"].SetFillStyle(1001)
            self.observable_settings["jetscape_distribution"].SetMarkerSize(0.0)
            self.observable_settings["jetscape_distribution"].SetMarkerStyle(0)
            self.observable_settings["jetscape_distribution"].SetLineWidth(0)
            self.observable_settings["jetscape_distribution"].DrawCopy("E3 same")
        elif self.observable_settings["jetscape_distribution"].GetNbinsX() == 1:
            self.observable_settings["jetscape_distribution"].SetMarkerSize(self.marker_size)
            self.observable_settings["jetscape_distribution"].SetMarkerStyle(self.data_marker + 1)
            self.observable_settings["jetscape_distribution"].SetMarkerColor(self.jetscape_color[0])
            self.observable_settings["jetscape_distribution"].SetLineStyle(self.line_style)
            self.observable_settings["jetscape_distribution"].SetLineWidth(self.line_width)
            self.observable_settings["jetscape_distribution"].SetLineColor(self.jetscape_color[0])
            self.observable_settings["jetscape_distribution"].DrawCopy("PE same")
        legend.AddEntry(self.observable_settings["jetscape_distribution"], model_display_name, "f")

        legend.Draw()

        # Draw ratio
        pad2.cd()
        if self.observable_settings["data_distribution"]:
            data_ratio = self.plot_utils.divide_tgraph_by_tgraph(
                self.observable_settings["data_distribution"], self.observable_settings["data_distribution"]
            )
            data_ratio.Draw("PE Z same")
            self.output_dict[f"data_ratio_{label}"] = data_ratio
        if self.observable_settings["ratio"]:
            self.output_dict[f"ratio_{label}"] = self.observable_settings["ratio"]
            if self.observable_settings["ratio"].GetN() > 1:
                self.observable_settings["ratio"].SetFillColor(self.jetscape_color[0])
                self.observable_settings["ratio"].SetFillColorAlpha(self.jetscape_color[0], self.jetscape_alpha[0])
                self.observable_settings["ratio"].SetFillStyle(1001)
                self.observable_settings["ratio"].SetMarkerSize(0.0)
                self.observable_settings["ratio"].SetMarkerStyle(0)
                self.observable_settings["ratio"].SetLineWidth(0)
                self.observable_settings["ratio"].Draw("E3 same")
            elif self.observable_settings["ratio"].GetN() == 1:
                self.observable_settings["ratio"].SetMarkerSize(self.marker_size)
                self.observable_settings["ratio"].SetMarkerStyle(self.data_marker + 1)
                self.observable_settings["ratio"].SetMarkerColor(self.jetscape_color[0])
                self.observable_settings["ratio"].SetLineStyle(self.line_style)
                self.observable_settings["ratio"].SetLineWidth(self.line_width)
                self.observable_settings["ratio"].SetLineColor(self.jetscape_color[0])
                self.observable_settings["ratio"].Draw("PE same")

        if self.observable_settings["ratio"]:
            legend_ratio.AddEntry(self.observable_settings["ratio"], f"{model_display_name}/Data", "f")
        if self.observable_settings["data_distribution"]:
            legend_ratio.AddEntry(data_ratio, "Data uncertainties", "PE")
        legend_ratio.Draw()

        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920 + 2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()

        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.25
        text_latex.SetTextSize(0.065)
        text = f"#bf{{{observable_type}_{observable}}} #sqrt{{#it{{s}}}} = {self.sqrts / 1000.0} TeV"
        text_latex.DrawLatex(x, 0.83, text)
        text = f"{centrality} {self.suffix} {pt_suffix}"
        text_latex.DrawLatex(x, 0.73, text)

        if getattr(self, "_area_normalized", False):
            text_latex.DrawLatex(x, 0.63, "MC area-normalized to data")

        if self.skip_pp:
            text = "skip data plot -- no pp data in HEPData"
            text_latex.DrawLatex(x, 0.53, text)

        pad2.cd()
        if self.skip_pp_ratio:
            text = "skip ratio plot -- pp/AA binning mismatch"
            text_latex.DrawLatex(x, 0.93, text)

        c.SaveAs(str(self.output_dir / f"{self.hname}{self.file_format}"))
        # Also emit a log-x/log-y twin of this distribution plot (NOT for R_AA ratios).
        self._save_logxy_twin(
            c, self.hname, dist_pad=pad1, ratio_pad=pad2, blank_histo=myBlankHisto, log_ymin=log_ymin, log_ymax=log_ymax
        )
        c.Close()

    # -------------------------------------------------------------------------------------------
    # Plot event QA
    # -------------------------------------------------------------------------------------------
    def plot_event_qa(self):
        h_pt_hat = self.input_file.Get("h_pt_hat")
        h_pt_hat_weighted = self.input_file.Get("h_pt_hat_weighted")
        if not h_pt_hat or not h_pt_hat_weighted:
            logger.warning("h_pt_hat not found, skipping event QA")
            return
        normalization_uncertainty = 0
        sum_weights_integral = 1  # avoid division by zero
        sum_weights = 0
        xsec = 1  # avoid division by zero
        xsec_error = 0
        if self.is_AA:
            for centrality in self.observable_centrality_list:
                # Only plot those centralities that exist
                # Retrieve the h_weight_sum histogram for the current centrality bin
                h_weight_sum = self.input_file.Get(f"h_weight_sum_{centrality}")
                # Check if the histogram exists and if its content is non-zero
                if not h_weight_sum or h_weight_sum.GetBinContent(1) == 0:
                    continue
                logger.info(centrality)

                # Crosscheck that pt-hat weighting is satisfactory
                # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
                sum_weights = self.input_file.Get(f"h_weight_sum_{centrality}").GetBinContent(1)
                h_pt_hat = self.input_file.Get("h_pt_hat")
                h_pt_hat_weighted = self.input_file.Get("h_pt_hat_weighted")

                # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
                h_pt_hat.Scale(1.0 / sum_weights, "width")
                h_pt_hat_weighted.Scale(1.0 / sum_weights, "width")

                # Compute normalization uncertainty (binned approximation)
                h_weights = self.input_file.Get(f"h_weights_{centrality}")
                sum_weights_integral = 0
                normalization_uncertainty = 0
                for i in range(1, h_weights.GetNbinsX() + 1):
                    sum_weights_integral += h_weights.GetBinCenter(i) * h_weights.GetBinContent(i)
                    normalization_uncertainty = np.sqrt(
                        np.square(normalization_uncertainty)
                        + h_weights.GetBinContent(i) * np.square(h_weights.GetBinCenter(i))
                    )
                logger.info(f"sum_weights {centrality}: {sum_weights}")
                logger.info(f"sum_weights_integral {centrality}: {sum_weights_integral}")
                logger.info(
                    f"normalization_uncertainty {centrality}: {100 * normalization_uncertainty / sum_weights_integral} %"
                )

                # Also compute overall pt-hat cross-section uncertainty
                h_xsec = self.input_file.Get(f"h_xsec_{centrality}")
                xsec = h_xsec.GetBinContent(1)
                xsec_error = self.input_file.Get(f"h_xsec_error_{centrality}").GetBinContent(1)
                logger.info(f"xsec {centrality}: {xsec / h_xsec.GetEntries()}")
                logger.info(f"xsec_uncertainty {centrality}: {100 * xsec_error / xsec}\n")

        else:
            # Crosscheck that pt-hat weighting is satisfactory
            # Make sure that the large-weight, low-pt-hat range has sufficient statistics to avoid normalization fluctuations
            sum_weights = self.input_file.Get("h_weight_sum").GetBinContent(1)
            h_pt_hat = self.input_file.Get("h_pt_hat")
            h_pt_hat_weighted = self.input_file.Get("h_pt_hat_weighted")

            # Normalize by sum of weights, i.e. 1/sigma_pt_hat * dsigma/dpt_hat
            h_pt_hat.Scale(1.0 / sum_weights, "width")
            h_pt_hat_weighted.Scale(1.0 / sum_weights, "width")

            # Compute normalization uncertainty (binned approximation)
            h_weights = self.input_file.Get("h_weights")
            sum_weights_integral = 0
            normalization_uncertainty = 0
            for i in range(1, h_weights.GetNbinsX() + 1):
                sum_weights_integral += h_weights.GetBinCenter(i) * h_weights.GetBinContent(i)
                normalization_uncertainty = np.sqrt(
                    np.square(normalization_uncertainty)
                    + h_weights.GetBinContent(i) * np.square(h_weights.GetBinCenter(i))
                )
            logger.info(f"sum_weights: {sum_weights}")
            logger.info(f"sum_weights_integral: {sum_weights_integral}")
            logger.info(f"normalization_uncertainty: {100 * normalization_uncertainty / sum_weights_integral} %")

            # Also compute overall pt-hat cross-section uncertainty
            h_xsec = self.input_file.Get("h_xsec")
            xsec = h_xsec.GetBinContent(1)
            xsec_error = self.input_file.Get("h_xsec_error").GetBinContent(1)
            logger.info(f"xsec: {xsec / h_xsec.GetEntries()}")
            logger.info(f"xsec_uncertainty: {100 * xsec_error / xsec}")

        c = ROOT.TCanvas("c", "c", 600, 650)
        c.Draw()
        c.cd()

        # Distribution
        pad2_dy = 0.45
        pad1 = ROOT.TPad("myPad", "The pad", 0, pad2_dy, 1, 1)
        pad1.SetLeftMargin(0.2)
        pad1.SetTopMargin(0.08)
        pad1.SetRightMargin(0.04)
        pad1.SetBottomMargin(0.0)
        pad1.SetTicks(0, 1)
        pad1.Draw()
        pad1.SetLogy()
        pad1.SetLogx()
        pad1.cd()

        legend = ROOT.TLegend(0.58, 0.65, 0.75, 0.9)
        self.plot_utils.setup_legend(legend, 0.055, sep=-0.1, title=f"#sqrt{{s_{{NN}}}} = {self.sqrts} GeV")

        self.bins = np.array(h_pt_hat.GetXaxis().GetXbins())
        myBlankHisto = ROOT.TH1F("myBlankHisto", "Blank Histogram", 1, self.bins[0], self.bins[-1])
        myBlankHisto.SetNdivisions(505)
        myBlankHisto.SetXTitle("#hat{p}_{T} (GeV/#it{c})")
        myBlankHisto.SetYTitle("#frac{1}{#sigma_{#hat{p}_{T}}} #frac{d#sigma}{d#hat{p}_{T}}")
        myBlankHisto.SetMaximum(np.power(10, 5 + (self.power - 3)))
        myBlankHisto.SetMinimum(2e-15)  # Don't draw 0 on top panel
        myBlankHisto.GetYaxis().SetTitleSize(0.08)
        myBlankHisto.GetYaxis().SetTitleOffset(1.1)
        myBlankHisto.GetYaxis().SetLabelSize(0.06)
        myBlankHisto.Draw("E")

        # Ratio
        c.cd()
        pad2 = ROOT.TPad("pad2", "pad2", 0, 0.02, 1, pad2_dy)
        pad2.SetTopMargin(0)
        pad2.SetBottomMargin(0.21)
        pad2.SetLeftMargin(0.2)
        pad2.SetRightMargin(0.04)
        pad2.SetTicks(0, 1)
        pad2.SetLogx()
        pad2.Draw()
        pad2.cd()

        myBlankHisto2 = myBlankHisto.Clone("myBlankHisto_C")
        myBlankHisto2.SetYTitle("#frac{weighted}{ideal}")
        myBlankHisto2.SetXTitle("#hat{p}_{T} (GeV/#it{c})")
        myBlankHisto2.GetXaxis().SetTitleSize(26)
        myBlankHisto2.GetXaxis().SetTitleFont(43)
        myBlankHisto2.GetXaxis().SetTitleOffset(2.3)
        myBlankHisto2.GetXaxis().SetLabelFont(43)
        myBlankHisto2.GetXaxis().SetLabelSize(22)
        myBlankHisto2.GetYaxis().SetTitleSize(28)
        myBlankHisto2.GetYaxis().SetTitleFont(43)
        myBlankHisto2.GetYaxis().SetTitleOffset(2.0)
        myBlankHisto2.GetYaxis().SetLabelFont(43)
        myBlankHisto2.GetYaxis().SetLabelSize(20)
        myBlankHisto2.GetYaxis().SetNdivisions(505)
        myBlankHisto2.GetYaxis().SetRangeUser(0.0, 1.99)
        myBlankHisto2.Draw("")

        # Draw generated pt-hat
        pad1.cd()
        h_pt_hat.SetMarkerSize(self.marker_size)
        h_pt_hat.SetMarkerStyle(self.data_marker)
        h_pt_hat.SetMarkerColor(self.data_color)
        h_pt_hat.SetLineStyle(self.line_style)
        h_pt_hat.SetLineWidth(self.line_width)
        h_pt_hat.SetLineColor(self.data_color)
        h_pt_hat.Draw("PE Z same")
        legend.AddEntry(h_pt_hat, f"#alpha={self.power}, generated", "PE")

        # Draw weighted pt-hat
        h_pt_hat_weighted.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted.SetLineStyle(self.line_style)
        h_pt_hat_weighted.SetLineWidth(self.line_width)
        h_pt_hat_weighted.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted.Draw("PE Z same")
        legend.AddEntry(h_pt_hat_weighted, f"#alpha={self.power}, weighted", "PE")

        # Draw generated pt-hat weighted by inverse of ideal weight (pt_hat/pt_ref)^alpha
        h_pt_hat_weighted_ideal = h_pt_hat.Clone()
        h_pt_hat_weighted_ideal.SetName("h_pt_hat_weighted_ideal")
        for i in range(1, h_pt_hat_weighted_ideal.GetNbinsX() + 1):
            content = h_pt_hat_weighted_ideal.GetBinContent(i)
            low_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinLowEdge(i)
            up_edge = h_pt_hat_weighted_ideal.GetXaxis().GetBinUpEdge(i)
            ideal_weight = (np.power(up_edge, self.power + 1) - np.power(low_edge, self.power + 1)) / (
                (up_edge - low_edge) * np.power(self.pt_ref, self.power) * (self.power + 1)
            )
            # ideal_weight = np.power( h_pt_hat_weighted_ideal.GetXaxis().GetBinCenter(i) / self.pt_ref, self.power)
            h_pt_hat_weighted_ideal.SetBinContent(i, content / ideal_weight)
            h_pt_hat_weighted_ideal.SetBinError(i, 0.0)
        h_pt_hat_weighted_ideal.SetLineColorAlpha(ROOT.kGreen - 8, self.jetscape_alpha[0])
        h_pt_hat_weighted_ideal.SetLineWidth(3)
        h_pt_hat_weighted_ideal.Draw("L same")
        legend.AddEntry(h_pt_hat_weighted_ideal, f"#alpha={self.power}, ideal weight", "L")

        legend.Draw()

        # Plot ratio of weighted to ideal weighted
        pad2.cd()
        h_pt_hat_weighted_ratio = h_pt_hat_weighted.Clone()
        h_pt_hat_weighted_ratio.SetName("h_pt_hat_weighted_ratio")
        h_pt_hat_weighted_ratio.Divide(h_pt_hat_weighted_ideal)
        h_pt_hat_weighted_ratio.SetMarkerSize(self.marker_size)
        h_pt_hat_weighted_ratio.SetMarkerStyle(self.data_marker)
        h_pt_hat_weighted_ratio.SetMarkerColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.SetLineStyle(self.line_style)
        h_pt_hat_weighted_ratio.SetLineWidth(self.line_width)
        h_pt_hat_weighted_ratio.SetLineColor(self.jetscape_color[0])
        h_pt_hat_weighted_ratio.Draw("PE Z same")

        line = ROOT.TLine(self.bins[0], 1, self.bins[-1], 1)
        line.SetLineColor(920 + 2)
        line.SetLineStyle(2)
        line.SetLineWidth(2)
        line.Draw()

        pad1.cd()
        text_latex = ROOT.TLatex()
        text_latex.SetNDC()

        x = 0.25
        text_latex.SetTextSize(0.06)
        text = f"#sigma_{{n_{{event}}}} = {100 * normalization_uncertainty / sum_weights_integral:.2f}%"
        text_latex.DrawLatex(x, 0.83, text)
        text = f"#sigma_{{xsec}} = {100 * xsec_error / xsec:.2f}%"
        text_latex.DrawLatex(x, 0.76, text)

        c.SaveAs(str(self.output_dir / f"pt_hat{self.file_format}"))
        c.Close()

    # ---------------------------------------------------------------
    # Save all ROOT histograms to file
    # ---------------------------------------------------------------
    def write_output_objects(self):
        logger.info("\nWriting output objects...")

        # Save histograms to ROOT file
        output_ROOT_filename = self.output_dir / "final_results.root"
        f_ROOT = ROOT.TFile(str(output_ROOT_filename), "recreate")
        f_ROOT.cd()
        for key, val in self.output_dict.items():
            if val:
                val.SetName(key)
                val.Write()
                if isinstance(val, (ROOT.TH1)):
                    val.SetDirectory(0)
                    del val
        f_ROOT.Close()

        # Also save JETSCAPE AA/pp ratios as numpy arrays to HDF5 file
        if self.is_AA:
            with uproot.open(output_ROOT_filename) as uproot_file:
                output_HDF5_filename = self.output_dir / "final_results.h5"
                with h5py.File(output_HDF5_filename, "w") as hf:
                    keys = [key.split(";", 1)[0] for key in uproot_file.keys()]  # noqa: SIM118
                    for key in keys:
                        if "jetscape" in key:
                            bin_edges = uproot_file[key].axis().edges()
                            values = uproot_file[key].values()
                            errors = uproot_file[key].errors()

                            hf.create_dataset(f"{key}_bin_edges", data=bin_edges)
                            hf.create_dataset(f"{key}_values", data=values)
                            hf.create_dataset(f"{key}_errors", data=errors)

    # -------------------------------------------------------------------------------------------
    # Generate pptx of one plot per slide, for convenience
    # -------------------------------------------------------------------------------------------
    def generate_pptx(self) -> None:
        # Delayed import since this isn't on available on many systems.
        # Provide a message for the user on how they could potentially proceed.
        try:
            import pptx  # pyright: ignore [reportMissingImports] # noqa: PLC0415
        except ImportError:
            _msg = """The `pptx` package is required to generate a presentation, but could not be imported.

If you have powerpoint available on your system, please install it separately with:

    pip install python-pptx

Once it's installed, please try again.
"""

        # Create a blank presentation
        p = pptx.Presentation()

        # Set slide layouts
        title_slide_layout = p.slide_layouts[0]
        blank_slide_layout = p.slide_layouts[6]

        # Make a title slide
        slide = p.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"QA for {self.sqrts / 1000.0} TeV"
        author = slide.placeholders[1]
        author.text = "STAT WG"

        # Loop through all output files and plot
        files = [f for f in self.output_dir.iterdir() if str(f).endswith(self.file_format)]
        for file in sorted(files):
            img = self.output_dir / file
            slide = p.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(
                str(img), left=pptx.util.Inches(2.0), top=pptx.util.Inches(1.0), width=pptx.util.Inches(5.0)
            )

        p.save(str(self.output_dir / "results.pptx"))


def main_entry_point() -> None:
    """Main entry point for plotting STAT results."""
    # Define arguments
    parser = argparse.ArgumentParser(description="Plot JETSCAPE events")
    parser.add_argument(
        "-c",
        "--configFile",
        action="store",
        type=Path,
        metavar="configFile",
        help="Path of config file for analysis (e.g. config/STAT_5020.yaml)",
    )
    parser.add_argument(
        "-i",
        "--inputFile",
        action="store",
        type=Path,
        metavar="inputFile",
        default=Path("pp_5020_plot/histograms_5020_merged.root"),
        help="Input file",
    )
    parser.add_argument(
        "-o",
        "--outputDir",
        action="store",
        type=Path,
        metavar="outputDir",
        default=None,
        help="Output directory for output to be written to",
    )
    parser.add_argument(
        "-r",
        "--refFile",
        action="store",
        type=Path,
        metavar="refFile",
        default=Path("final_results.root"),
        help="pp reference file",
    )
    parser.add_argument(
        "-l",
        "--logLevel",
        action="store",
        type=str,
        metavar="logLevel",
        default="INFO",
        help="Set the logging level (DEBUG, INFO, WARNING, ERROR, CRITICAL)",
    )
    parser.add_argument(
        "-m",
        "--model-name",
        action="store",
        type=str,
        required=False,
        metavar="model_name",
        default="",
        # NOTE(RJE): Autodetection is not straightforward in this case since we've lost the model name from the filename.
        #            So we just opt for backwards compatibility.
        help="Name of the model which we are analyzing. Default: '', which indicates that we are using jetscape.",
    )

    # Parse the arguments
    args = parser.parse_args()

    # Setup logging
    helpers.setup_logging(level=args.logLevel)

    logger.info("Executing plot_results_STAT...\n")

    # Validation and setup
    if not args.configFile.exists():
        msg = f'File "{args.configFile}" does not exist! Exiting!'
        raise ValueError(msg)
    if not args.inputFile.exists():
        msg = f'File "{args.inputFile}" does not exist! Exiting!'
        raise ValueError(msg)

    analysis = PlotResults(
        config_file=args.configFile, input_file=args.inputFile, output_dir=args.outputDir, pp_ref_file=args.refFile
    )
    try:
        analysis.plot_results()
    except Exception as e:
        logger.exception(e)
        sys.exit(1)


if __name__ == "__main__":
    main_entry_point()
